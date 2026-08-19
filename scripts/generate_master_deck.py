"""Integrated Master Deck Generator (PPTX + PDF) for International Basketball Analytics (2005-2024).

Generates both:
1. presentation/International_Basketball_Analytics_Presentation.pptx (Editable widescreen 16:9)
2. presentation/International_Basketball_Analytics_Presentation.pdf (Widescreen 16:9 960x540 pt)

MVP-34 Final Presentation Polish:
- Slide 10: SHRINKAGE_VISUAL (Bayesian shrinkage flow + conceptual slider)
- Slide 14: MONTE_CARLO_VISUAL (Stylized tournament bracket + 3 large simulation metrics)
- Slide 15: WORKSPACE_MOCKUP (Analyst decision UI mockup with anti-hindsight quarantine)
- Slide 27: ROLE_VALUE_TREE (Club organizational value tree: Coach / Scout / Director)
- Slide 28: ROADMAP_TIMELINE (Horizontal 30-day club integration roadmap)
- All other 25 slides remain 100% intact and canonical.
"""

import os
import sys
from pathlib import Path

# PPTX imports
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# ReportLab imports
from reportlab.pdfgen import canvas
from reportlab.lib import colors

PROJECT_ROOT = Path("f:/España2005-2025")
PRESENTATION_DIR = PROJECT_ROOT / "presentation"
OUTPUT_PPTX = PRESENTATION_DIR / "International_Basketball_Analytics_Presentation.pptx"
OUTPUT_PDF = PRESENTATION_DIR / "International_Basketball_Analytics_Presentation.pdf"

# ---------------------------------------------------------------------------
# COLOR PALETTE TOKENS
# ---------------------------------------------------------------------------
DARK_BG = RGBColor(10, 17, 40)        # #0A1128 Deep Navy
DARK_CARD = RGBColor(18, 30, 66)      # #121E42
DARK_TEXT = RGBColor(248, 250, 252)   # #F8FAFC
DARK_MUTED = RGBColor(148, 163, 184)  # #94A3B8

LIGHT_BG = RGBColor(248, 250, 252)    # #F8FAFC
CARD_BG = RGBColor(255, 255, 255)     # #FFFFFF
CARD_BORDER = RGBColor(226, 232, 240) # #E2E8F0
TEXT_MAIN = RGBColor(15, 23, 42)      # #0F172A
TEXT_MUTED = RGBColor(71, 85, 105)    # #475569

CYAN_ACCENT = RGBColor(2, 132, 199)   # #0284C7 Primary Cyan
ORANGE_ACCENT = RGBColor(234, 88, 12) # #EA580C Basketball Orange
GREEN_ACCENT = RGBColor(22, 163, 74)  # #16A34A Success Green
AMBER_ACCENT = RGBColor(217, 119, 6)  # #D97706 Amber
PURPLE_ACCENT = RGBColor(124, 58, 237)# #7C3AED Purple

# ReportLab Color equivalents
RL_DARK_BG = colors.HexColor("#0A1128")
RL_DARK_CARD = colors.HexColor("#121E42")
RL_DARK_TEXT = colors.HexColor("#F8FAFC")
RL_DARK_MUTED = colors.HexColor("#94A3B8")

RL_LIGHT_BG = colors.HexColor("#F8FAFC")
RL_CARD_BG = colors.HexColor("#FFFFFF")
RL_CARD_BORDER = colors.HexColor("#E2E8F0")
RL_TEXT_MAIN = colors.HexColor("#0F172A")
RL_TEXT_MUTED = colors.HexColor("#475569")

RL_CYAN = colors.HexColor("#0284C7")
RL_ORANGE = colors.HexColor("#EA580C")
RL_GREEN = colors.HexColor("#16A34A")
RL_AMBER = colors.HexColor("#D97706")
RL_PURPLE = colors.HexColor("#7C3AED")


# ---------------------------------------------------------------------------
# SLIDE DEFINITIONS (30 SLIDES, BALANCED ACROSS 6 ACTS)
# ---------------------------------------------------------------------------
SLIDES_DATA = [
    # SLIDE 1: COVER
    {
        "type": "COVER",
        "tag": "PORTFOLIO EXECUTIVE DECK",
        "title": "International Basketball Analytics",
        "subtitle": "20 Years of FIBA Competitions (2005-2024): An End-to-End Data, ML & Decision Support System",
        "meta": "Python + R + DuckDB + Parquet  |  1,145 Games  |  27,353 Player Records  |  227 Automated Tests",
        "author": "Miguel - Basketball Data Analyst & Quantitative Scout"
    },
    # SLIDE 2: EXECUTIVE SUMMARY
    {
        "type": "EXEC_SUMMARY",
        "category": "PROJECT AT A GLANCE",
        "title": "Executive Summary: From Raw Data to Court Decision",
        "subtitle": "A complete data engineering and predictive pipeline built for professional basketball decision-makers.",
        "cards": [
            ("01. Data Foundation", "18 senior FIBA tournaments (EuroBasket, World Cup, Olympics) structured into an OLAP DuckDB warehouse with 12 normalized tables.", CYAN_ACCENT),
            ("02. Advanced Metrics", "Dean Oliver Four Factors normalized to pace, True Shooting %, and longitudinal career stability via Bayesian shrinkage.", ORANGE_ACCENT),
            ("03. Machine Learning", "Chronological 17-fold walk-forward validation (1,105 test games). Brier Score 0.1967 (vs 0.2500 naive) and ECE 0.0314.", GREEN_ACCENT),
            ("04. Decision Support", "1.5-page pre-game coaching briefs, 6 objective player archetypes, and anti-hindsight validation mode.", PURPLE_ACCENT)
        ],
        "footer": "<b>Core Value</b>: Reducing tactical uncertainty while preserving coaching intuition."
    },
    # SLIDE 3: THE PROBLEM IN BASKETBALL
    {
        "type": "SPLIT_PROBLEM",
        "category": "ACT I: THE PROBLEM",
        "title": "The Analytics Challenge in Professional Basketball",
        "subtitle": "Why traditional boxscores and 40-page scouting dossiers fail modern coaching staffs.",
        "left_title": "The Status Quo: Information Overload",
        "left_items": [
            "Scouting reports are flooded with raw totals (PPG, RPG) that ignore pace and lineup context.",
            "Short-tournament variance (6-9 games) misleads decision-makers into over-valuing shooting hot streaks.",
            "Models trained with random k-fold suffer from temporal data leakage, learning future tactical eras.",
            "Analysts deliver disconnected charts instead of structured pre-game decision questions."
        ],
        "right_title": "Our Solution: The Data-First Pipeline",
        "right_items": [
            "Pace-neutral Four Factors isolate true possession efficiency.",
            "Bayesian shrinkage (lambda = 0.75) stabilizes small-sample noise toward historical priors.",
            "Strict walk-forward expanding window guarantees zero future contamination.",
            "Concise 1.5-page coaching briefs delivering actionable tactical hypotheses before tip-off."
        ]
    },
    # SLIDE 4: SECTION DIVIDER 1
    {
        "type": "SECTION_DIVIDER",
        "number": "01",
        "title": "Data Architecture & Engineering",
        "subtitle": "From Raw Historical Archives to In-Process Columnar OLAP with DuckDB and Parquet"
    },
    # SLIDE 5: DATA SCALE & PROVENANCE
    {
        "type": "BIG_METRICS",
        "category": "DATA FOUNDATION",
        "title": "Historical Scope & Canonical Scale",
        "subtitle": "Two decades of international men's senior basketball verified with SHA-256 cryptographic provenance.",
        "metrics": [
            ("1,145", "OFFICIAL MATCHES", "Every game across 18 FIBA tournaments (2005-2024)."),
            ("18", "TOURNAMENTS", "EuroBasket (2005-2022), FIBA World Cup (2006-2023), Olympics (2008-2024)."),
            ("27,353", "PLAYER PERFORMANCES", "Individual boxscore game logs across all 18 tournaments."),
            ("2,124", "CANONICAL PLAYERS", "Deterministic entity resolution mapping names across 20 years.")
        ],
        "footer": "<b>Integrity</b>: 100% immutable data lake with zero external live API dependencies."
    },
    # SLIDE 6: MEDALLION ARCHITECTURE
    {
        "type": "FLOW_ARCHITECTURE",
        "category": "DATA ENGINEERING",
        "title": "Medallion Data Warehouse Architecture",
        "subtitle": "Decoupled pipeline separating ingestion, relational validation, columnar storage, and analytics.",
        "steps": [
            ("01_RAW", "HTML/JSON Lake", "18 tournament source archives stored for full offline reproducibility.", CYAN_ACCENT),
            ("02_VALIDATED", "DuckDB Warehouse", "12 relational tables (28.5 MB uncompressed) with 3NF integrity.", ORANGE_ACCENT),
            ("03_ANALYTICS", "Apache Parquet", "11 high-performance columnar marts optimized for vectorized queries.", GREEN_ACCENT),
            ("04_DELIVERY", "Python & R APIs", "Dual-engine consumption for ML training, statistical EDA, and briefs.", PURPLE_ACCENT)
        ],
        "footer": "<b>Performance</b>: Sub-millisecond aggregations on 27k+ player rows directly in-memory."
    },
    # SLIDE 7: QA & MATHEMATICAL INVARIANTS
    {
        "type": "QA_INVARIANTS",
        "category": "DATA QUALITY",
        "title": "Automated Quality Assurance & Invariants",
        "subtitle": "Rigorous mathematical contracts enforced prior to any statistical modeling.",
        "cards": [
            ("200 Minutes / Game Rule", "Every 40-minute regulation game strictly verifies exactly 200 player-minutes (plus 25 min per 5-min overtime). Any mismatch fails the pipeline.", CYAN_ACCENT),
            ("Score & Boxscore Closure", "Sum of points (2P + 3P + FT) is reconciled against official final match scores. Zero scoring discrepancies across 1,145 games.", GREEN_ACCENT),
            ("Entity Resolution Engine", "Unified 2,124 players across language variants, accents, and name changes (e.g. Dirk Nowitzki, Bojan vs Bogdan Bogdanovic) using custom slug hashing.", ORANGE_ACCENT)
        ],
        "footer": "<b>Quality Standard</b>: Corrupt data yields corrupt decisions. We validate before we model."
    },
    # SLIDE 8: SECTION DIVIDER 2
    {
        "type": "SECTION_DIVIDER",
        "number": "02",
        "title": "Analytics & Modeling Engine",
        "subtitle": "Pace Normalization, Longitudinal Stability, Player Archetypes, and Calibrated Machine Learning"
    },
    # SLIDE 9: FOUR FACTORS & PACE
    {
        "type": "FOUR_FACTORS",
        "category": "METRICS FRAMEWORK",
        "title": "Pace Normalization & Dean Oliver Four Factors",
        "subtitle": "Decomposing basketball efficiency into four fundamental, pace-neutral drivers.",
        "factors": [
            ("Effective FG% (eFG%)", "Shooting Efficiency", "Weights 3-pointers at 1.5x. Accounts for shot quality differences.", "eFG% = (FGM + 0.5 * 3PM) / FGA", CYAN_ACCENT),
            ("Turnover Rate (TOV%)", "Ball Security", "Measures possessions ending in turnover per 100 plays.", "TOV% = TOV / (FGA + 0.44 * FTA + TOV)", ORANGE_ACCENT),
            ("Off. Rebounding (ORB%)", "Second Chances", "Percentage of available offensive rebounds secured.", "ORB% = ORB / (ORB + Opp_DRB)", GREEN_ACCENT),
            ("Free Throw Rate (FTR)", "Foul Drawing", "Ratio of free throw attempts to field goal attempts.", "FTR = FTA / FGA", PURPLE_ACCENT)
        ],
        "footer": "<b>Pace Neutrality</b>: Comparing 65-possession games directly against 85-possession games."
    },
    # SLIDE 10: POLISHED - PLAYER ANALYTICS & SHRINKAGE VISUAL
    {
        "type": "SHRINKAGE_VISUAL",
        "category": "LONGITUDINAL INFERENCE",
        "title": "Overcoming Small-Sample Tournament Noise",
        "subtitle": "Bayesian shrinkage (lambda = 0.75) stabilizes noisy tournament shooting toward career priors.",
        "flow_steps": [
            ("1. RAW TOURNAMENT SIGNAL", "7-game sample: 15/30 3PT = 50.0% (Small-Sample Trap)", ORANGE_ACCENT),
            ("2. HIGH VARIANCE & NOISE", "A swing of just 3 shots drops 3PT% from 50% to 40%", ORANGE_ACCENT),
            ("3. BAYESIAN SHRINKAGE (lambda = 0.75)", "Weighted contraction toward historical career baseline", CYAN_ACCENT),
            ("4. STABILIZED ESTIMATE", "Reliable true skill projection for pre-game tactical planning", GREEN_ACCENT)
        ],
        "spectrum_items": [
            ("Raw Tournament Estimate (50.0%)", "High Variance / Noise", ORANGE_ACCENT),
            ("Shrunk Stabilized Estimate (38.5%)", "Reliable Game-Plan Prior", CYAN_ACCENT),
            ("Global Population Mean (34.2%)", "Historical Baseline", GREEN_ACCENT)
        ],
        "callout": ("lambda = 0.75", "BAYESIAN SHRINKAGE FACTOR", "Contracts short-sample variance toward 3,767 player career baselines."),
        "bootstrap": "Cluster Bootstrap (B = 5,000): Non-parametric 95% confidence intervals across 3,767 qualified campaigns (>=40 min).",
        "footer": "<b>Core Principle</b>: Small samples exaggerate hot streaks. Shrinkage restores empirical reality."
    },
    # SLIDE 11: 6 OBJECTIVE FUNCTIONAL ARCHETYPES
    {
        "type": "ARCHETYPES_GRID",
        "category": "PLAYER TAXONOMY",
        "title": "6 Objective Functional Archetypes",
        "subtitle": "Clustering 3,767 player campaigns via K-Means++ and PCA on 14 normalized per-40 metrics.",
        "archetypes": [
            ("Primary Initiator", "High usage, primary P&R handler, elite shot creation and passing volume.", CYAN_ACCENT),
            ("Floor Spacer", "High 3P volume, catch-and-shoot specialist, low interior touches, elite spacing.", ORANGE_ACCENT),
            ("Interior Hub", "Post playmaker, rim finisher, high offensive rebounding and short-roll passer.", GREEN_ACCENT),
            ("Floor General", "Pure playmaker, exceptional assist-to-turnover ratio, low turnover rate, tempo controller.", PURPLE_ACCENT),
            ("Defensive Anchor", "Rim protector, elite defensive rebounder, screen-and-roll finisher, low usage.", CYAN_ACCENT),
            ("Balanced Wing", "Versatile two-way wing, secondary creator, switchable defender, balanced shot profile.", ORANGE_ACCENT)
        ],
        "footer": "<b>Scouting Value</b>: Moving beyond rigid 1-5 positions to evaluate true roster complement."
    },
    # SLIDE 12: MACHINE LEARNING CALIBRATION
    {
        "type": "ML_CALIBRATION",
        "category": "PREDICTIVE MODELING",
        "title": "Calibrated Machine Learning: LightGBM",
        "subtitle": "Probabilistic match forecasting regularized with L2 penalty and isotonic calibration.",
        "metrics": [
            ("0.1967", "BRIER SCORE", "vs 0.2500 naive baseline (+21.3% quadratic error reduction)."),
            ("0.0314", "CALIBRATION ERROR (ECE)", "Expected Calibration Error: within 3.14% of perfect calibration."),
            ("11.74", "POINT MARGIN MAE", "Mean Absolute Error in point spread across 1,105 test matches.")
        ],
        "footer": "<b>Honest Probability</b>: An event assigned 70% probability occurs exactly 7 out of 10 times."
    },
    # SLIDE 13: 17-FOLD WALK-FORWARD
    {
        "type": "WALK_FORWARD",
        "category": "VALIDATION METHODOLOGY",
        "title": "17-Fold Temporal Walk-Forward Protocol",
        "subtitle": "Strict expanding-window splits preventing temporal data leakage.",
        "timeline": [
            ("Fold 01", "Train: EuroBasket 2005", "Test: FIBA World Cup 2006 (Out-of-sample)"),
            ("Fold 05", "Train: 2005-2009 Tournaments", "Test: FIBA World Cup 2010 (Out-of-sample)"),
            ("Fold 11", "Train: 2005-2016 Tournaments", "Test: FIBA EuroBasket 2017 (Out-of-sample)"),
            ("Fold 17", "Train: 2005-2023 Tournaments", "Test: Paris Olympic Games 2024 (Out-of-sample)")
        ],
        "footer": "<b>Why Walk-Forward Matters</b>: Random K-Fold leaks modern 3-point volume into past eras."
    },
    # SLIDE 14: POLISHED - TOURNAMENT SIMULATION ENGINE
    {
        "type": "MONTE_CARLO_VISUAL",
        "category": "TOURNAMENT PROJECTIONS",
        "title": "Monte Carlo Tournament Simulation Engine",
        "subtitle": "Simulating full tournament outcome distributions across 180,000 bracket iterations.",
        "pipeline": "Historical Data  ->  Strength Priors  ->  Monte Carlo (180,000)  ->  Probabilistic Outcomes",
        "bracket_rounds": [
            ("Round of 16", ["Match 01", "Match 02", "Match 03", "Match 04"]),
            ("Quarterfinals", ["QF 01", "QF 02"]),
            ("Semifinals", ["SF 01", "SF 02"]),
            ("Final & Podium", ["CHAMPION", "MEDAL TREE"])
        ],
        "metrics": [
            ("180,000", "BRACKET ITERATIONS", "Full dynamic bracket simulations accounting for match sequence and cumulative fatigue.", CYAN_ACCENT),
            ("lambda = 0.75", "BAYESIAN SHRINKAGE", "Prevents over-reacting to preliminary round group-stage upsets.", ORANGE_ACCENT),
            ("PROBABILISTIC", "OUTCOME DISTRIBUTIONS", "Outputs conditional podium probability trees rather than fragile single-winner picks.", GREEN_ACCENT)
        ],
        "footer": "<b>Modeling Focus</b>: The system simulates outcome distributions, not deterministic winners."
    },
    # SLIDE 15: POLISHED - ANALYST WORKSPACE & ANTI-HINDSIGHT
    {
        "type": "WORKSPACE_MOCKUP",
        "category": "DECISION SUPPORT WORKSPACE",
        "title": "Interactive Workspace & Anti-Hindsight Protocol",
        "subtitle": "Isolating pre-game evidence from post-game outcomes to eliminate hindsight bias.",
        "app_title": "ANALYST DECISION WORKSPACE  |  Streamlit Decision Support Engine",
        "sidebar_items": [
            "Target Match: Spain vs USA",
            "Checkpoint: T-1 (Pre-Game)",
            "Evidence: Four Factors & Drop",
            "Roster State: Canonical 2008",
            "Status: ANTI-HINDSIGHT ACTIVE"
        ],
        "pre_game_items": [
            "Pace-neutral Four Factors (+4.2 Net Rating half-court)",
            "Opponent defense alert: Center deep drop on P&R",
            "Longitudinal shooting priors (lambda = 0.75 stabilized)",
            "Functional archetype matchup synergy & rim delta"
        ],
        "post_game_locked": [
            "[QUARANTINED POST-GAME INFORMATION]",
            "- Final score & official boxscore totals",
            "- Shooting percentages of current game",
            "- Post-hoc tactical rationalizations"
        ],
        "timeline_steps": [
            ("T-7 Pre-Brief", CYAN_ACCENT),
            ("T-1 Decision Log", ORANGE_ACCENT),
            ("Game Tip-Off", GREEN_ACCENT),
            ("Outcome Revealed", PURPLE_ACCENT),
            ("Post-Game Audit", CYAN_ACCENT)
        ],
        "footer": "<b>Anti-Hindsight Guarantee</b>: Analyst decisions are permanently timestamped before tip-off."
    },
    # SLIDE 16: SECTION DIVIDER 3
    {
        "type": "SECTION_DIVIDER",
        "number": "03",
        "title": "Selected Case Studies",
        "subtitle": "Four Applied Demonstrations: Tactical Support, Data Engineering, ML Rigor, and Longitudinal R Analysis"
    },
    # SLIDE 17: CASE STUDY 1 - TACTICAL
    {
        "type": "CASE_STUDY",
        "category": "CASE STUDY 01 : TACTICAL DECISION SUPPORT",
        "title": "From Data to the Whiteboard: Beijing 2008 Final",
        "subtitle": "How pre-game analytics framed the tactical strategy for Spain vs. USA 'Redeem Team'.",
        "left_title": "Quantitative Pre-Game Signals",
        "left_items": [
            "Context: Spain lost by 37 pts (82-119) in group phase against USA.",
            "Half-Court Signal: In 5v5 half-court sets, Spain held a +4.2 Net Rating advantage via Gasol brothers inside.",
            "Transition Vulnerability: USA scored 1.25 PPP on transition off turnovers.",
            "Defensive Flaw: USA bigs dropped deep into paint to pack driving lanes."
        ],
        "right_title": "Tactical Whiteboard Brief",
        "right_items": [
            "Pace Control: Limit total possessions below 72 using 16+ second attacks.",
            "Transition Defense: Implement 2-3 matchup zone immediately after made baskets.",
            "Pick-and-Pop Spacing: Exploit deep drop with 3-pointers from Pau, Marc, and Garbajosa.",
            "Court Outcome: Spain cut deficit to 4 pts (108-104) at 2:20; final 107-118."
        ],
        "footer": "<b>Demonstration</b>: Turning complex data into 3 concrete tactical questions for the head coach."
    },
    # SLIDE 18: CASE STUDY 2 - DATA ENGINEERING
    {
        "type": "CASE_STUDY",
        "category": "CASE STUDY 02 : SPORTS DATA ENGINEERING",
        "title": "In-Process OLAP Lakehouse with DuckDB & Parquet",
        "subtitle": "Deterministic data engineering for 20 years of heterogeneous international tournaments.",
        "left_title": "Engineering Challenges",
        "left_items": [
            "Heterogeneous FIBA Boxscores across 2005-2024 with varying column conventions.",
            "Diacritic Name Variations across international alphabets (Nowitzki, Spanoulis, Bogdanovic).",
            "Zero Cloud Infrastructure cost constraint: 100% offline local reproducibility.",
            "Dual-Language Access: Zero-copy querying from both Python and R."
        ],
        "right_title": "Architecture Solutions",
        "right_items": [
            "DuckDB Relational Engine: 12 tables, Star Schema, in-process ACID transactions.",
            "11 Apache Parquet Marts: Snappy-compressed columnar storage with SHA-256 hashes.",
            "Deterministic Entity Resolver: 2,124 canonical player IDs without paid commercial APIs.",
            "Automated QA Suite: 227 tests in pytest enforcing 200 min/game invariants."
        ],
        "footer": "<b>Demonstration</b>: Production-quality data engineering capable of scaling to domestic leagues."
    },
    # SLIDE 19: CASE STUDY 3 - ML WALK-FORWARD
    {
        "type": "CASE_STUDY",
        "category": "CASE STUDY 03 : MACHINE LEARNING RIGOR",
        "title": "Zero-Leakage ML with 17 Chronological Folds",
        "subtitle": "Avoiding temporal leakage and delivering calibrated probabilities for tournament play.",
        "left_title": "Methodological Traps Avoided",
        "left_items": [
            "Random K-Fold Cross-Validation leaks future playstyle evolution into past evaluations.",
            "Over-confident predictions: raw neural nets produce uncalibrated probabilities in short events.",
            "Overfitting on small samples: complex models over-index on single-game tournament anomalies.",
            "Causal Overclaiming: confusing feature importance with tactical ground truth."
        ],
        "right_title": "Rigorous ML Protocol",
        "right_items": [
            "Expanding Window: 17 folds evaluating strictly unseen tournaments (1,105 test games).",
            "LightGBM with L2 Regularization on pre-game Four Factors differentials.",
            "Isotonic Probability Calibration: Brier Score 0.1967 (vs 0.2500 naive), ECE 0.0314.",
            "TreeSHAP Attribution: Interpretable local feature contribution without causal exaggeration."
        ],
        "footer": "<b>Demonstration</b>: Scientific honesty and calibrated probability estimation."
    },
    # SLIDE 20: CASE STUDY 4 - LONGITUDINAL INFERENCE
    {
        "type": "CASE_STUDY",
        "category": "CASE STUDY 04 : R & STATISTICAL EDA",
        "title": "Longitudinal Shooting & Role Mining in R / Quarto",
        "subtitle": "Multi-year player evaluation, non-parametric bootstrap, and reproducible reporting.",
        "left_title": "Statistical Analysis in R",
        "left_items": [
            "Direct DuckDB Connection via DBI: Zero data duplication between Python and R.",
            "Longitudinal Career Tracking: 3,767 qualified player campaigns (>=40 minutes).",
            "Cluster Bootstrap (B = 5,000): Non-parametric 95% confidence intervals for True Shooting %.",
            "Rebounding & Turnover Stability: Higher year-over-year correlation (r > 0.65) than 3P%."
        ],
        "right_title": "Clustering & Quarto Reporting",
        "right_items": [
            "PCA Decomposition: First 3 principal components explain >60% of player variance.",
            "K-Means++ Clustering: Discovery of 6 functional archetypes based on 14 normalized metrics.",
            "Quarto CLI Pipeline: Automated compilation into interactive HTML report with ggplot2.",
            "Custom Visualization: `theme_basketball_analytics()` with 300 DPI vector exports."
        ],
        "footer": "<b>Demonstration</b>: Advanced R, tidyverse, statistical inference, and reproducible reporting."
    },
    # SLIDE 21: SECTION DIVIDER 4
    {
        "type": "SECTION_DIVIDER",
        "number": "04",
        "title": "Validation & Engineering Rigor",
        "subtitle": "227 Automated Pytest Tests, Cross-Language Verification, and Deterministic Master Runner"
    },
    # SLIDE 22: 227 AUTOMATED TESTS
    {
        "type": "TESTING_SUITE",
        "category": "ENGINEERING QUALITY",
        "title": "227 Automated Regression Tests (100% Pass)",
        "subtitle": "Continuous integration test suite covering 26 specialized test modules.",
        "cards": [
            ("Data Integrity & Invariants", "Validates 200 min/game, boxscore closure, SHA-256 hashes, DuckDB schemas, and column contracts across all 18 tournaments.", CYAN_ACCENT),
            ("Analytical & ML Precision", "Tests walk-forward temporal splits, Brier Score bounds, ECE limits, Monte Carlo distribution bounds, and SHAP attribution coherence.", ORANGE_ACCENT),
            ("Portfolio & Cross-Language", "Verifies R-DuckDB parity, Quarto compilation, publication package completeness, and zero prohibited marketing buzzwords.", GREEN_ACCENT)
        ],
        "footer": "<b>Execution Speed</b>: Full suite executes in ~2 minutes on standard hardware."
    },
    # SLIDE 23: CROSS-LANGUAGE PARITY
    {
        "type": "CROSS_LANGUAGE",
        "category": "INTEROPERABILITY",
        "title": "Cross-Language Parity: Python + R + DuckDB",
        "subtitle": "Ensuring identical analytical results regardless of execution environment.",
        "items": [
            ("Unified Warehouse", "Single `basketball_analytics.duckdb` file accessed concurrently by Python (`duckdb`) and R (`DBI::dbConnect`).", CYAN_ACCENT),
            ("Zero Discrepancy", "Exact match on 1,145 games, 2,290 team rows, 27,353 player boxscores, and 4,350 qualified records.", GREEN_ACCENT),
            ("Metric Synchronization", "Identical eFG% (0.5355), Pace (61.72), and 3P Attempt Rate (0.3153) across Python and R pipelines.", ORANGE_ACCENT),
            ("Decoupled Design", "Python handles ETL/ML; R handles statistical inference and visualization; Parquet provides columnar handoff.", PURPLE_ACCENT)
        ],
        "footer": "<b>Verified</b>: Deterministic script `scripts/verify_cross_language.py` checks parity automatically."
    },
    # SLIDE 24: REPRODUCIBILITY IN 1 COMMAND
    {
        "type": "REPRODUCIBILITY",
        "category": "ONE-COMMAND EXECUTION",
        "title": "Deterministic End-to-End Reproducibility",
        "subtitle": "Run the entire project verification pipeline with a single master command.",
        "code": "python scripts/run_project.py",
        "steps": [
            ("1. Environment Check", "Verifies Python 3.10+, R 4.4+, DuckDB, PyArrow, LightGBM, Pandas, and Pytest.", CYAN_ACCENT),
            ("2. Warehouse Check", "Validates DuckDB database (16.8 MB) and 11 Parquet analytical marts.", ORANGE_ACCENT),
            ("3. R Analysis Pipeline", "Executes R statistical scripts and generates figures in reports/figures_r/ in ~21s.", GREEN_ACCENT),
            ("4. Pytest Regression Suite", "Runs all 227 automated tests with 100% pass rate in ~2 minutes.", PURPLE_ACCENT)
        ],
        "footer": "<b>Zero Live Dependencies</b>: Everything needed is packaged inside the repository."
    },
    # SLIDE 25: SECTION DIVIDER 5
    {
        "type": "SECTION_DIVIDER",
        "number": "05",
        "title": "Limitations & Professional Scope",
        "subtitle": "Scientific Honesty, Transparent Boundaries, and Operational Value for Basketball Clubs"
    },
    # SLIDE 26: WHAT WE CAN VS CANNOT CLAIM
    {
        "type": "LIMITATIONS_MATRIX",
        "category": "METHODOLOGICAL BOUNDARIES",
        "title": "Transparent Boundaries: What the System Claims",
        "subtitle": "Clear distinction between data-backed capabilities and operational limits.",
        "left_title": "What the System DEMONSTRATES",
        "left_items": [
            "Pace-neutral Four Factors and True Shooting evaluation.",
            "Longitudinal Bayesian shrinkage to control short-tournament variance.",
            "Strict out-of-sample calibrated match probability modeling.",
            "Structured pre-game briefs that ask the right tactical questions.",
            "Fully reproducible in-process OLAP data lakehouse."
        ],
        "right_title": "What the System DOES NOT Claim",
        "right_items": [
            "No 25Hz live optical tracking (no Second Spectrum camera telemetry).",
            "No automatic game decisions: data structures evidence, coach decides.",
            "No infallible betting prediction: models estimate uncertainty.",
            "No causal guarantees: statistical associations are not mechanical laws.",
            "No real-time in-game bench tracking (focused on pre/post game)."
        ],
        "footer": "<b>Scientific Integrity</b>: Honest analytics earns the trust of coaching staffs."
    },
    # SLIDE 27: POLISHED - VALUE BY BASKETBALL ROLE
    {
        "type": "ROLE_VALUE_TREE",
        "category": "ORGANIZATIONAL IMPACT",
        "title": "Operational Value by Basketball Role",
        "subtitle": "How different stakeholders in a basketball club extract value from this platform.",
        "club_header": "BASKETBALL CLUB DECISION ECOSYSTEM",
        "roles": [
            ("HEAD COACH & STAFF", "WHAT TO CHANGE", [
                "Pre-game 1.5-page coaching briefs",
                "Opponent P&R drop coverage alerts",
                "Pace & possession targets (<= 72 poss)"
            ], CYAN_ACCENT),
            ("SCOUTING DEPARTMENT", "WHO FITS", [
                "6 objective functional archetypes",
                "Small-sample shooting shrinkage",
                "Roster complement & role audit"
            ], ORANGE_ACCENT),
            ("SPORTING DIRECTOR / GM", "WHAT TO INVEST IN", [
                "Evidence-based contract valuation",
                "Career trajectory stability curves",
                "Recruitment risk minimization"
            ], GREEN_ACCENT)
        ],
        "value_flow": "RAW DATA (DuckDB Lakehouse)  ->  EVIDENCE (R & Calibrated ML)  ->  DECISION SUPPORT (Pizarra & Roster)",
        "footer": "<b>Operational Impact</b>: Moving from raw data tables directly to court and roster decisions."
    },
    # SLIDE 28: POLISHED - DAY-1 CLUB INTEGRATION PLAN
    {
        "type": "ROADMAP_TIMELINE",
        "category": "PRACTICAL ONBOARDING",
        "title": "Day-1 Club Onboarding Roadmap (30 Days)",
        "subtitle": "How this modular platform deploys seamlessly in domestic club leagues.",
        "timeline_header": "DAY 01 --------------> DAY 10 --------------> DAY 20 --------------> DAY 30",
        "phases": [
            ("PHASE 01 : DATA AUDIT", "Days 1-10", [
                "Connect domestic boxscores (Synergy, Genius, EuroLeague API)",
                "Validate relational DuckDB schema & table contracts",
                "Execute automated QA suite (200 min/game invariants)"
            ], CYAN_ACCENT),
            ("PHASE 02 : ANALYST WORKFLOW", "Days 11-20", [
                "Calibrate league-specific Four Factors baselines",
                "Set Bayesian shrinkage priors (lambda = 0.75) for roster",
                "Train video analysts on anti-hindsight workspace"
            ], ORANGE_ACCENT),
            ("PHASE 03 : STAFF INTEGRATION", "Days 21-30", [
                "Deliver automated pre-game briefs directly to coaching staff",
                "Integrate Four Factors alerts into whiteboard pre-game talks",
                "Establish weekly post-game accuracy review loop"
            ], GREEN_ACCENT)
        ],
        "milestone": "DAY 30 MILESTONE: FULLY EMBEDDED DECISION SYSTEM  (Zero Cloud Overhead, 100% Offline Reproducibility)",
        "footer": "<b>Turnkey Deployment</b>: Modular pipeline ready for immediate adoption in professional clubs."
    },
    # SLIDE 29: KEY TAKEAWAYS
    {
        "type": "TAKEAWAYS",
        "category": "CONCLUSION",
        "title": "Summary & Key Takeaways",
        "subtitle": "A disciplined, reproducible approach to basketball quantitative analytics.",
        "points": [
            ("1. Data-First & Rigorous", "20 years of FIBA tournaments structured in DuckDB and Parquet with 227 automated tests and zero data leakage.", CYAN_ACCENT),
            ("2. Calibrated Machine Learning", "Brier Score 0.1967 and ECE 0.0314 across 17 chronological folds (1,105 test matches).", ORANGE_ACCENT),
            ("3. Actionable Tactical Impact", "1.5-page briefs that ask concrete questions for the whiteboard instead of dumping raw statistical charts.", GREEN_ACCENT),
            ("4. Transparent & Reproducible", "Complete open-source pipeline executed with a single command on GitHub.", PURPLE_ACCENT)
        ],
        "footer": "<b>Core Philosophy</b>: Data structures evidence. The coach decides the game."
    },
    # SLIDE 30: REPOSITORY & CONTACT
    {
        "type": "CLOSING",
        "tag": "PROJECT REPOSITORY & CITATION",
        "title": "International Basketball Analytics",
        "subtitle": "Explore the Full Open-Source Codebase, Case Studies, and Automated Tests",
        "github_url": "https://github.com/miguelo0203/basketball-analytics",
        "release": "Release v1.0.0 (Data-First, Verified, MIT License)",
        "citation": "Citation: CITATION.cff (CFF v1.2.0)  |  Author: Miguel"
    }
]


# ===========================================================================
# POWERPOINT (.PPTX) GENERATOR
# ===========================================================================
def build_pptx():
    """Generate the complete redesigned 30-slide PPTX deck."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    for idx, data in enumerate(SLIDES_DATA):
        slide = prs.slides.add_slide(blank_layout)
        stype = data["type"]

        if stype == "COVER":
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = DARK_BG
            bg.line.fill.background()

            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(3.2), Inches(0.4))
            badge.fill.solid()
            badge.fill.fore_color.rgb = CYAN_ACCENT
            badge.line.fill.background()
            tf_b = badge.text_frame
            p_b = tf_b.paragraphs[0]
            p_b.text = data["tag"]
            p_b.font.size = Pt(11)
            p_b.font.bold = True
            p_b.font.color.rgb = DARK_TEXT

            tbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.8))
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(40)
            p.font.bold = True
            p.font.color.rgb = DARK_TEXT

            sbox = slide.shapes.add_textbox(Inches(1.0), Inches(3.4), Inches(11.3), Inches(1.0))
            tf_s = sbox.text_frame
            tf_s.word_wrap = True
            p_s = tf_s.paragraphs[0]
            p_s.text = data["subtitle"]
            p_s.font.size = Pt(18)
            p_s.font.color.rgb = DARK_MUTED

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.8), Inches(11.333), Inches(1.5))
            card.fill.solid()
            card.fill.fore_color.rgb = DARK_CARD
            card.line.fill.background()
            tf_c = card.text_frame
            tf_c.word_wrap = True
            tf_c.margin_left = Inches(0.4)
            tf_c.margin_top = Inches(0.3)
            p_c1 = tf_c.paragraphs[0]
            p_c1.text = data["meta"]
            p_c1.font.size = Pt(14)
            p_c1.font.bold = True
            p_c1.font.color.rgb = CYAN_ACCENT

            p_c2 = tf_c.add_paragraph()
            p_c2.text = data["author"]
            p_c2.font.size = Pt(13)
            p_c2.font.color.rgb = DARK_TEXT

        elif stype == "SECTION_DIVIDER":
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = DARK_BG
            bg.line.fill.background()

            num_box = slide.shapes.add_textbox(Inches(1.2), Inches(1.4), Inches(3.0), Inches(1.6))
            tf_n = num_box.text_frame
            p_n = tf_n.paragraphs[0]
            p_n.text = data["number"]
            p_n.font.size = Pt(72)
            p_n.font.bold = True
            p_n.font.color.rgb = CYAN_ACCENT

            line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(1.2), Inches(3.2), Inches(2.5), Inches(0.08))
            line.fill.solid()
            line.fill.fore_color.rgb = ORANGE_ACCENT
            line.line.fill.background()

            tbox = slide.shapes.add_textbox(Inches(1.2), Inches(3.5), Inches(10.8), Inches(1.4))
            tf = tbox.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = DARK_TEXT

            sbox = slide.shapes.add_textbox(Inches(1.2), Inches(5.0), Inches(10.8), Inches(1.2))
            tf_s = sbox.text_frame
            tf_s.word_wrap = True
            p_s = tf_s.paragraphs[0]
            p_s.text = data["subtitle"]
            p_s.font.size = Pt(18)
            p_s.font.color.rgb = DARK_MUTED

        elif stype == "CLOSING":
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = DARK_BG
            bg.line.fill.background()

            badge = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(1.0), Inches(3.8), Inches(0.4))
            badge.fill.solid()
            badge.fill.fore_color.rgb = GREEN_ACCENT
            badge.line.fill.background()
            p_b = badge.text_frame.paragraphs[0]
            p_b.text = data["tag"]
            p_b.font.size = Pt(11)
            p_b.font.bold = True
            p_b.font.color.rgb = DARK_TEXT

            tbox = slide.shapes.add_textbox(Inches(1.0), Inches(1.6), Inches(11.3), Inches(1.4))
            p = tbox.text_frame.paragraphs[0]
            p.text = data["title"]
            p.font.size = Pt(36)
            p.font.bold = True
            p.font.color.rgb = DARK_TEXT

            sbox = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(11.3), Inches(0.8))
            p_s = sbox.text_frame.paragraphs[0]
            p_s.text = data["subtitle"]
            p_s.font.size = Pt(18)
            p_s.font.color.rgb = DARK_MUTED

            card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(4.2), Inches(11.333), Inches(2.2))
            card.fill.solid()
            card.fill.fore_color.rgb = DARK_CARD
            card.line.fill.background()
            tf_c = card.text_frame
            tf_c.margin_left = Inches(0.5)
            tf_c.margin_top = Inches(0.4)

            p1 = tf_c.paragraphs[0]
            p1.text = "GitHub Repository: " + data["github_url"]
            p1.font.size = Pt(16)
            p1.font.bold = True
            p1.font.color.rgb = CYAN_ACCENT

            p2 = tf_c.add_paragraph()
            p2.text = data["release"]
            p2.font.size = Pt(14)
            p2.font.color.rgb = GREEN_ACCENT

            p3 = tf_c.add_paragraph()
            p3.text = data["citation"]
            p3.font.size = Pt(13)
            p3.font.color.rgb = DARK_TEXT

        else:
            # Standard Content Slide
            bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(13.333), Inches(7.5))
            bg.fill.solid()
            bg.fill.fore_color.rgb = LIGHT_BG
            bg.line.fill.background()

            cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.3))
            p_cat = cat_box.text_frame.paragraphs[0]
            p_cat.text = data.get("category", "INTERNATIONAL BASKETBALL ANALYTICS").upper()
            p_cat.font.size = Pt(10)
            p_cat.font.bold = True
            p_cat.font.color.rgb = CYAN_ACCENT

            tbox = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.6))
            p_t = tbox.text_frame.paragraphs[0]
            p_t.text = data["title"]
            p_t.font.size = Pt(22)
            p_t.font.bold = True
            p_t.font.color.rgb = TEXT_MAIN

            sbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.3), Inches(11.7), Inches(0.5))
            p_s = sbox.text_frame.paragraphs[0]
            p_s.text = data.get("subtitle", "")
            p_s.font.size = Pt(13)
            p_s.font.color.rgb = TEXT_MUTED

            # Layout Dispatcher
            if stype == "EXEC_SUMMARY":
                cards = data["cards"]
                coords = [
                    (Inches(0.8), Inches(1.9), Inches(5.6), Inches(2.2)),
                    (Inches(6.8), Inches(1.9), Inches(5.6), Inches(2.2)),
                    (Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.2)),
                    (Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.2))
                ]
                for (title, desc, color), (left, top, w, h) in zip(cards, coords):
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
                    p1 = tf.paragraphs[0]
                    p1.text = title
                    p1.font.size = Pt(15)
                    p1.font.bold = True
                    p1.font.color.rgb = color
                    p2 = tf.add_paragraph()
                    p2.text = desc
                    p2.font.size = Pt(12)
                    p2.font.color.rgb = TEXT_MAIN

            elif stype == "BIG_METRICS":
                metrics = data["metrics"]
                w = Inches(2.7)
                for i, (val, label, sub) in enumerate(metrics):
                    left = Inches(0.8 + i * 2.95)
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), w, Inches(4.4))
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = Inches(0.25)
                    tf.margin_top = Inches(0.4)

                    p1 = tf.paragraphs[0]
                    p1.text = val
                    p1.font.size = Pt(36)
                    p1.font.bold = True
                    p1.font.color.rgb = CYAN_ACCENT

                    p2 = tf.add_paragraph()
                    p2.text = label
                    p2.font.size = Pt(12)
                    p2.font.bold = True
                    p2.font.color.rgb = TEXT_MAIN

                    p3 = tf.add_paragraph()
                    p3.text = sub
                    p3.font.size = Pt(11)
                    p3.font.color.rgb = TEXT_MUTED

            elif stype in ["SPLIT_PROBLEM", "CASE_STUDY", "LIMITATIONS_MATRIX"]:
                left_c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.7), Inches(4.5))
                left_c.fill.solid()
                left_c.fill.fore_color.rgb = CARD_BG
                left_c.line.color.rgb = CARD_BORDER
                tf_l = left_c.text_frame
                tf_l.word_wrap = True
                tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = Inches(0.3)
                p_lt = tf_l.paragraphs[0]
                p_lt.text = data.get("left_title", "Problem Analysis")
                p_lt.font.size = Pt(16)
                p_lt.font.bold = True
                p_lt.font.color.rgb = ORANGE_ACCENT

                for item in data.get("left_items", []):
                    p = tf_l.add_paragraph()
                    p.text = "• " + item
                    p.font.size = Pt(11.5)
                    p.font.color.rgb = TEXT_MAIN

                right_c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.9), Inches(5.7), Inches(4.5))
                right_c.fill.solid()
                right_c.fill.fore_color.rgb = CARD_BG
                right_c.line.color.rgb = CARD_BORDER
                tf_r = right_c.text_frame
                tf_r.word_wrap = True
                tf_r.margin_left = tf_r.margin_right = tf_r.margin_top = Inches(0.3)
                p_rt = tf_r.paragraphs[0]
                p_rt.text = data.get("right_title", "Tactical Solution")
                p_rt.font.size = Pt(16)
                p_rt.font.bold = True
                p_rt.font.color.rgb = CYAN_ACCENT

                for item in data.get("right_items", []):
                    p = tf_r.add_paragraph()
                    p.text = "• " + item
                    p.font.size = Pt(11.5)
                    p.font.color.rgb = TEXT_MAIN

            elif stype == "FLOW_ARCHITECTURE":
                steps = data["steps"]
                w = Inches(2.7)
                for i, (layer, name, desc, color) in enumerate(steps):
                    left = Inches(0.8 + i * 2.95)
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), w, Inches(4.4))
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = Inches(0.25)
                    tf.margin_top = Inches(0.35)

                    p1 = tf.paragraphs[0]
                    p1.text = layer
                    p1.font.size = Pt(13)
                    p1.font.bold = True
                    p1.font.color.rgb = color

                    p2 = tf.add_paragraph()
                    p2.text = name
                    p2.font.size = Pt(18)
                    p2.font.bold = True
                    p2.font.color.rgb = TEXT_MAIN

                    p3 = tf.add_paragraph()
                    p3.text = desc
                    p3.font.size = Pt(11.5)
                    p3.font.color.rgb = TEXT_MUTED

            elif stype == "ARCHETYPES_GRID":
                archetypes = data["archetypes"]
                coords = [
                    (Inches(0.8), Inches(1.9), Inches(3.7), Inches(2.1)),
                    (Inches(4.8), Inches(1.9), Inches(3.7), Inches(2.1)),
                    (Inches(8.8), Inches(1.9), Inches(3.7), Inches(2.1)),
                    (Inches(0.8), Inches(4.3), Inches(3.7), Inches(2.1)),
                    (Inches(4.8), Inches(4.3), Inches(3.7), Inches(2.1)),
                    (Inches(8.8), Inches(4.3), Inches(3.7), Inches(2.1))
                ]
                for (name, desc, color), (left, top, w, h) in zip(archetypes, coords):
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.2)
                    p1 = tf.paragraphs[0]
                    p1.text = name
                    p1.font.size = Pt(15)
                    p1.font.bold = True
                    p1.font.color.rgb = color
                    p2 = tf.add_paragraph()
                    p2.text = desc
                    p2.font.size = Pt(11)
                    p2.font.color.rgb = TEXT_MAIN

            elif stype == "FOUR_FACTORS":
                factors = data["factors"]
                coords = [
                    (Inches(0.8), Inches(1.9), Inches(5.6), Inches(2.1)),
                    (Inches(6.8), Inches(1.9), Inches(5.6), Inches(2.1)),
                    (Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.1)),
                    (Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.1))
                ]
                for (name, subtitle, desc, formula, color), (left, top, w, h) in zip(factors, coords):
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.2)
                    p1 = tf.paragraphs[0]
                    p1.text = name + " — " + subtitle
                    p1.font.size = Pt(14)
                    p1.font.bold = True
                    p1.font.color.rgb = color
                    p2 = tf.add_paragraph()
                    p2.text = desc
                    p2.font.size = Pt(11)
                    p2.font.color.rgb = TEXT_MAIN
                    p3 = tf.add_paragraph()
                    p3.text = formula
                    p3.font.size = Pt(10)
                    p3.font.bold = True
                    p3.font.color.rgb = CYAN_ACCENT

            elif stype == "ML_CALIBRATION":
                metrics = data["metrics"]
                for i, (val, label, sub) in enumerate(metrics):
                    left = Inches(0.8 + i * 3.9)
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.0), Inches(3.7), Inches(4.4))
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = Inches(0.3)
                    tf.margin_top = Inches(0.5)

                    p1 = tf.paragraphs[0]
                    p1.text = val
                    p1.font.size = Pt(36)
                    p1.font.bold = True
                    p1.font.color.rgb = CYAN_ACCENT

                    p2 = tf.add_paragraph()
                    p2.text = label
                    p2.font.size = Pt(14)
                    p2.font.bold = True
                    p2.font.color.rgb = TEXT_MAIN

                    p3 = tf.add_paragraph()
                    p3.text = sub
                    p3.font.size = Pt(12)
                    p3.font.color.rgb = TEXT_MUTED

            elif stype == "WALK_FORWARD":
                timeline = data["timeline"]
                for i, (fold, train, test) in enumerate(timeline):
                    top = Inches(1.9 + i * 1.15)
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.0))
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.3)
                    tf.margin_top = Inches(0.2)
                    p1 = tf.paragraphs[0]
                    p1.text = fold + " | " + train + " -> " + test
                    p1.font.size = Pt(13)
                    p1.font.bold = True
                    p1.font.color.rgb = TEXT_MAIN

            # MVP-34 POLISHED SLIDE 10: SHRINKAGE_VISUAL
            elif stype == "SHRINKAGE_VISUAL":
                # Left column: 4-stage flow
                left_c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.9), Inches(5.6), Inches(4.5))
                left_c.fill.solid()
                left_c.fill.fore_color.rgb = CARD_BG
                left_c.line.color.rgb = CARD_BORDER
                tf_l = left_c.text_frame
                tf_l.word_wrap = True
                tf_l.margin_left = tf_l.margin_right = tf_l.margin_top = Inches(0.25)

                p_lt = tf_l.paragraphs[0]
                p_lt.text = "ESTIMATION FLOW: FROM NOISE TO STABILITY"
                p_lt.font.size = Pt(13)
                p_lt.font.bold = True
                p_lt.font.color.rgb = CYAN_ACCENT

                for i, (step_t, step_d, col) in enumerate(data["flow_steps"]):
                    p_s1 = tf_l.add_paragraph()
                    p_s1.text = "▼  " + step_t if i > 0 else "●  " + step_t
                    p_s1.font.size = Pt(11.5)
                    p_s1.font.bold = True
                    p_s1.font.color.rgb = col

                    p_s2 = tf_l.add_paragraph()
                    p_s2.text = "    " + step_d
                    p_s2.font.size = Pt(10.5)
                    p_s2.font.color.rgb = TEXT_MUTED

                # Right column: Callout + Spectrum Slider + Bootstrap
                # 1. Top Callout Card
                call_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.9), Inches(5.7), Inches(1.6))
                call_card.fill.solid()
                call_card.fill.fore_color.rgb = DARK_CARD
                call_card.line.color.rgb = CYAN_ACCENT
                tf_cc = call_card.text_frame
                tf_cc.margin_left = Inches(0.3)
                tf_cc.margin_top = Inches(0.2)
                c_val, c_lbl, c_desc = data["callout"]
                p_cc1 = tf_cc.paragraphs[0]
                p_cc1.text = c_val + "  |  " + c_lbl
                p_cc1.font.size = Pt(15)
                p_cc1.font.bold = True
                p_cc1.font.color.rgb = CYAN_ACCENT
                p_cc2 = tf_cc.add_paragraph()
                p_cc2.text = c_desc
                p_cc2.font.size = Pt(11)
                p_cc2.font.color.rgb = DARK_TEXT

                # 2. Middle Spectrum Diagram Card
                spec_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(3.65), Inches(5.7), Inches(1.75))
                spec_card.fill.solid()
                spec_card.fill.fore_color.rgb = CARD_BG
                spec_card.line.color.rgb = CARD_BORDER
                tf_sp = spec_card.text_frame
                tf_sp.margin_left = Inches(0.3)
                tf_sp.margin_top = Inches(0.2)
                p_sp0 = tf_sp.paragraphs[0]
                p_sp0.text = "CONCEPTUAL SHRINKAGE SPECTRUM"
                p_sp0.font.size = Pt(12)
                p_sp0.font.bold = True
                p_sp0.font.color.rgb = TEXT_MAIN
                for label, desc, col in data["spectrum_items"]:
                    p_item = tf_sp.add_paragraph()
                    p_item.text = f"• {label} -> {desc}"
                    p_item.font.size = Pt(10.5)
                    p_item.font.bold = True
                    p_item.font.color.rgb = col

                # 3. Bottom Bootstrap Badge Card
                boot_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(5.55), Inches(5.7), Inches(0.85))
                boot_card.fill.solid()
                boot_card.fill.fore_color.rgb = CARD_BG
                boot_card.line.color.rgb = CARD_BORDER
                tf_bc = boot_card.text_frame
                tf_bc.margin_left = Inches(0.3)
                tf_bc.margin_top = Inches(0.15)
                p_bc = tf_bc.paragraphs[0]
                p_bc.text = data["bootstrap"]
                p_bc.font.size = Pt(10.5)
                p_bc.font.bold = True
                p_bc.font.color.rgb = GREEN_ACCENT

            # MVP-34 POLISHED SLIDE 14: MONTE_CARLO_VISUAL
            elif stype == "MONTE_CARLO_VISUAL":
                # Top pipeline banner
                pipe_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.55))
                pipe_box.fill.solid()
                pipe_box.fill.fore_color.rgb = DARK_CARD
                pipe_box.line.color.rgb = CYAN_ACCENT
                tf_p = pipe_box.text_frame
                tf_p.margin_left = Inches(0.3)
                p_p = tf_p.paragraphs[0]
                p_p.text = "SIMULATION PIPELINE: " + data["pipeline"]
                p_p.font.size = Pt(12)
                p_p.font.bold = True
                p_p.font.color.rgb = CYAN_ACCENT

                # Left side: Bracket representation
                br_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.5), Inches(5.6), Inches(3.9))
                br_card.fill.solid()
                br_card.fill.fore_color.rgb = CARD_BG
                br_card.line.color.rgb = CARD_BORDER
                tf_br = br_card.text_frame
                tf_br.margin_left = Inches(0.3)
                tf_br.margin_top = Inches(0.2)
                p_bt = tf_br.paragraphs[0]
                p_bt.text = "DYNAMIC TOURNAMENT BRACKET FLOW"
                p_bt.font.size = Pt(13)
                p_bt.font.bold = True
                p_bt.font.color.rgb = TEXT_MAIN

                for round_name, nodes in data["bracket_rounds"]:
                    p_r1 = tf_br.add_paragraph()
                    p_r1.text = "▶  " + round_name
                    p_r1.font.size = Pt(11)
                    p_r1.font.bold = True
                    p_r1.font.color.rgb = CYAN_ACCENT
                    p_r2 = tf_br.add_paragraph()
                    p_r2.text = "    " + "  |  ".join(nodes)
                    p_r2.font.size = Pt(10)
                    p_r2.font.color.rgb = TEXT_MUTED

                # Right side: 3 Large simulation metric cards
                for i, (val, lbl, desc, col) in enumerate(data["metrics"]):
                    top_y = Inches(2.5 + i * 1.33)
                    mc = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.8), top_y, Inches(5.7), Inches(1.22))
                    mc.fill.solid()
                    mc.fill.fore_color.rgb = CARD_BG
                    mc.line.color.rgb = CARD_BORDER
                    tf_m = mc.text_frame
                    tf_m.margin_left = Inches(0.3)
                    tf_m.margin_top = Inches(0.18)

                    p_m1 = tf_m.paragraphs[0]
                    p_m1.text = val + "  |  " + lbl
                    p_m1.font.size = Pt(14)
                    p_m1.font.bold = True
                    p_m1.font.color.rgb = col

                    p_m2 = tf_m.add_paragraph()
                    p_m2.text = desc
                    p_m2.font.size = Pt(10.5)
                    p_m2.font.color.rgb = TEXT_MAIN

            # MVP-34 POLISHED SLIDE 15: WORKSPACE_MOCKUP
            elif stype == "WORKSPACE_MOCKUP":
                # App Header
                app_hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.55))
                app_hdr.fill.solid()
                app_hdr.fill.fore_color.rgb = DARK_CARD
                app_hdr.line.color.rgb = CYAN_ACCENT
                tf_ah = app_hdr.text_frame
                tf_ah.margin_left = Inches(0.3)
                p_ah = tf_ah.paragraphs[0]
                p_ah.text = "💻  " + data["app_title"]
                p_ah.font.size = Pt(12)
                p_ah.font.bold = True
                p_ah.font.color.rgb = CYAN_ACCENT

                # Sidebar Panel
                side = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.5), Inches(3.2), Inches(3.9))
                side.fill.solid()
                side.fill.fore_color.rgb = DARK_BG
                side.line.color.rgb = DARK_CARD
                tf_sd = side.text_frame
                tf_sd.margin_left = Inches(0.25)
                tf_sd.margin_top = Inches(0.2)
                p_sdt = tf_sd.paragraphs[0]
                p_sdt.text = "APP FILTERS"
                p_sdt.font.size = Pt(11)
                p_sdt.font.bold = True
                p_sdt.font.color.rgb = CYAN_ACCENT
                for item in data["sidebar_items"]:
                    p_item = tf_sd.add_paragraph()
                    p_item.text = "• " + item
                    p_item.font.size = Pt(10)
                    p_item.font.color.rgb = DARK_TEXT

                # Main Evidence Panel: Left (Pre-game available) vs Right (Locked Post-game)
                pre_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(2.5), Inches(4.3), Inches(2.8))
                pre_card.fill.solid()
                pre_card.fill.fore_color.rgb = CARD_BG
                pre_card.line.color.rgb = GREEN_ACCENT
                tf_pr = pre_card.text_frame
                tf_pr.margin_left = Inches(0.25)
                tf_pr.margin_top = Inches(0.2)
                p_prt = tf_pr.paragraphs[0]
                p_prt.text = "PRE-GAME EVIDENCE (AVAILABLE)"
                p_prt.font.size = Pt(12)
                p_prt.font.bold = True
                p_prt.font.color.rgb = GREEN_ACCENT
                for item in data["pre_game_items"]:
                    p_pi = tf_pr.add_paragraph()
                    p_pi.text = "• " + item
                    p_pi.font.size = Pt(10)
                    p_pi.font.color.rgb = TEXT_MAIN

                post_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(8.7), Inches(2.5), Inches(3.8), Inches(2.8))
                post_card.fill.solid()
                post_card.fill.fore_color.rgb = CARD_BG
                post_card.line.color.rgb = ORANGE_ACCENT
                tf_po = post_card.text_frame
                tf_po.margin_left = Inches(0.25)
                tf_po.margin_top = Inches(0.2)
                p_pot = tf_po.paragraphs[0]
                p_pot.text = "POST-GAME (QUARANTINED)"
                p_pot.font.size = Pt(12)
                p_pot.font.bold = True
                p_pot.font.color.rgb = ORANGE_ACCENT
                for item in data["post_game_locked"]:
                    p_poi = tf_po.add_paragraph()
                    p_poi.text = item
                    p_poi.font.size = Pt(10)
                    p_poi.font.color.rgb = TEXT_MUTED

                # Bottom Timeline Ribbon
                rib_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(4.2), Inches(5.45), Inches(8.3), Inches(0.95))
                rib_card.fill.solid()
                rib_card.fill.fore_color.rgb = DARK_CARD
                rib_card.line.color.rgb = CYAN_ACCENT
                tf_rc = rib_card.text_frame
                tf_rc.margin_left = Inches(0.3)
                tf_rc.margin_top = Inches(0.18)
                p_rc = tf_rc.paragraphs[0]
                timeline_str = "  ->  ".join([f"{name}" for name, _ in data["timeline_steps"]])
                p_rc.text = "ANTI-HINDSIGHT AUDIT TRAIL:  " + timeline_str
                p_rc.font.size = Pt(11)
                p_rc.font.bold = True
                p_rc.font.color.rgb = DARK_TEXT

            # MVP-34 POLISHED SLIDE 27: ROLE_VALUE_TREE
            elif stype == "ROLE_VALUE_TREE":
                # Top Tree Header
                tree_hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.55))
                tree_hdr.fill.solid()
                tree_hdr.fill.fore_color.rgb = DARK_CARD
                tree_hdr.line.color.rgb = CYAN_ACCENT
                tf_th = tree_hdr.text_frame
                tf_th.margin_left = Inches(0.3)
                p_th = tf_th.paragraphs[0]
                p_th.text = "ORGANIZATIONAL IMPACT:  " + data["club_header"]
                p_th.font.size = Pt(12)
                p_th.font.bold = True
                p_th.font.color.rgb = CYAN_ACCENT

                # 3 Role Columns
                roles = data["roles"]
                w = Inches(3.7)
                for i, (role_name, role_q, items, col) in enumerate(roles):
                    left = Inches(0.8 + i * 3.9)
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), w, Inches(3.0))
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.margin_left = tf.margin_right = Inches(0.25)
                    tf.margin_top = Inches(0.25)

                    p1 = tf.paragraphs[0]
                    p1.text = role_name
                    p1.font.size = Pt(13)
                    p1.font.bold = True
                    p1.font.color.rgb = col

                    p2 = tf.add_paragraph()
                    p2.text = f'"{role_q}"'
                    p2.font.size = Pt(11)
                    p2.font.bold = True
                    p2.font.color.rgb = TEXT_MAIN

                    for item in items:
                        p_i = tf.add_paragraph()
                        p_i.text = "• " + item
                        p_i.font.size = Pt(10.5)
                        p_i.font.color.rgb = TEXT_MUTED

                # Bottom Value Flow Ribbon
                v_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.65), Inches(11.7), Inches(0.75))
                v_card.fill.solid()
                v_card.fill.fore_color.rgb = DARK_CARD
                v_card.line.color.rgb = GREEN_ACCENT
                tf_vc = v_card.text_frame
                tf_vc.margin_left = Inches(0.3)
                p_vc = tf_vc.paragraphs[0]
                p_vc.text = "VALUE PIPELINE:  " + data["value_flow"]
                p_vc.font.size = Pt(11)
                p_vc.font.bold = True
                p_vc.font.color.rgb = DARK_TEXT

            # MVP-34 POLISHED SLIDE 28: ROADMAP_TIMELINE
            elif stype == "ROADMAP_TIMELINE":
                # Top Timeline Progress Ribbon
                t_hdr = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.55))
                t_hdr.fill.solid()
                t_hdr.fill.fore_color.rgb = DARK_CARD
                t_hdr.line.color.rgb = CYAN_ACCENT
                tf_t = t_hdr.text_frame
                tf_t.margin_left = Inches(0.3)
                p_t = tf_t.paragraphs[0]
                p_t.text = "TRANSITION TIMELINE:  " + data["timeline_header"]
                p_t.font.size = Pt(12)
                p_t.font.bold = True
                p_t.font.color.rgb = CYAN_ACCENT

                # 3 Horizontal Phase Cards
                phases = data["phases"]
                w = Inches(3.7)
                for i, (p_name, p_days, p_items, col) in enumerate(phases):
                    left = Inches(0.8 + i * 3.9)
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.5), w, Inches(3.0))
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.margin_left = tf.margin_right = Inches(0.25)
                    tf.margin_top = Inches(0.25)

                    p1 = tf.paragraphs[0]
                    p1.text = p_name
                    p1.font.size = Pt(13)
                    p1.font.bold = True
                    p1.font.color.rgb = col

                    p2 = tf.add_paragraph()
                    p2.text = p_days
                    p2.font.size = Pt(11)
                    p2.font.bold = True
                    p2.font.color.rgb = TEXT_MAIN

                    for item in p_items:
                        p_i = tf.add_paragraph()
                        p_i.text = "• " + item
                        p_i.font.size = Pt(10.5)
                        p_i.font.color.rgb = TEXT_MUTED

                # Bottom Milestone Card
                m_card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(5.65), Inches(11.7), Inches(0.75))
                m_card.fill.solid()
                m_card.fill.fore_color.rgb = DARK_CARD
                m_card.line.color.rgb = GREEN_ACCENT
                tf_mc = m_card.text_frame
                tf_mc.margin_left = Inches(0.3)
                p_mc = tf_mc.paragraphs[0]
                p_mc.text = "⭐  " + data["milestone"]
                p_mc.font.size = Pt(11)
                p_mc.font.bold = True
                p_mc.font.color.rgb = GREEN_ACCENT

            elif stype == "CROSS_LANGUAGE":
                items = data["items"]
                for i, (title, desc, color) in enumerate(items):
                    top = Inches(1.9 + i * 1.15)
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), top, Inches(11.7), Inches(1.0))
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = Inches(0.3)
                    tf.margin_top = Inches(0.18)
                    p1 = tf.paragraphs[0]
                    p1.text = title + " — " + desc
                    p1.font.size = Pt(13)
                    p1.font.bold = True
                    p1.font.color.rgb = color

            elif stype == "REPRODUCIBILITY":
                cmd_box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(1.8), Inches(11.7), Inches(0.7))
                cmd_box.fill.solid()
                cmd_box.fill.fore_color.rgb = DARK_CARD
                cmd_box.line.color.rgb = CYAN_ACCENT
                tf_cmd = cmd_box.text_frame
                tf_cmd.margin_left = Inches(0.3)
                p_cmd = tf_cmd.paragraphs[0]
                p_cmd.text = "$ " + data["code"]
                p_cmd.font.size = Pt(14)
                p_cmd.font.bold = True
                p_cmd.font.color.rgb = CYAN_ACCENT

                steps = data["steps"]
                w = Inches(2.7)
                for i, (step_name, step_desc, color) in enumerate(steps):
                    left = Inches(0.8 + i * 2.95)
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(2.7), w, Inches(3.7))
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = Inches(0.2)
                    tf.margin_top = Inches(0.3)
                    p1 = tf.paragraphs[0]
                    p1.text = step_name
                    p1.font.size = Pt(14)
                    p1.font.bold = True
                    p1.font.color.rgb = color
                    p2 = tf.add_paragraph()
                    p2.text = step_desc
                    p2.font.size = Pt(11)
                    p2.font.color.rgb = TEXT_MAIN

            elif stype == "TAKEAWAYS":
                points = data["points"]
                coords = [
                    (Inches(0.8), Inches(1.9), Inches(5.6), Inches(2.2)),
                    (Inches(6.8), Inches(1.9), Inches(5.6), Inches(2.2)),
                    (Inches(0.8), Inches(4.3), Inches(5.6), Inches(2.2)),
                    (Inches(6.8), Inches(4.3), Inches(5.6), Inches(2.2))
                ]
                for (title, desc, color), (left, top, w, h) in zip(points, coords):
                    c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
                    c.fill.solid()
                    c.fill.fore_color.rgb = CARD_BG
                    c.line.color.rgb = CARD_BORDER
                    tf = c.text_frame
                    tf.word_wrap = True
                    tf.margin_left = tf.margin_right = tf.margin_top = Inches(0.25)
                    p1 = tf.paragraphs[0]
                    p1.text = title
                    p1.font.size = Pt(15)
                    p1.font.bold = True
                    p1.font.color.rgb = color
                    p2 = tf.add_paragraph()
                    p2.text = desc
                    p2.font.size = Pt(12)
                    p2.font.color.rgb = TEXT_MAIN

            else:
                cards = data.get("cards", data.get("roles", data.get("phases", [])))
                if cards:
                    w = Inches(3.7)
                    for i, item in enumerate(cards):
                        left = Inches(0.8 + i * 3.9)
                        c = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, Inches(1.9), w, Inches(4.5))
                        c.fill.solid()
                        c.fill.fore_color.rgb = CARD_BG
                        c.line.color.rgb = CARD_BORDER
                        tf = c.text_frame
                        tf.word_wrap = True
                        tf.margin_left = tf.margin_right = Inches(0.3)
                        tf.margin_top = Inches(0.4)

                        title = item[0]
                        desc = item[1]
                        color = item[2] if len(item) > 2 else CYAN_ACCENT

                        p1 = tf.paragraphs[0]
                        p1.text = title
                        p1.font.size = Pt(16)
                        p1.font.bold = True
                        p1.font.color.rgb = color if isinstance(color, RGBColor) else CYAN_ACCENT

                        p2 = tf.add_paragraph()
                        p2.text = desc
                        p2.font.size = Pt(12)
                        p2.font.color.rgb = TEXT_MAIN

            # Footer
            if "footer" in data:
                fbox = slide.shapes.add_textbox(Inches(0.8), Inches(6.6), Inches(11.7), Inches(0.4))
                p_f = fbox.text_frame.paragraphs[0]
                p_f.text = data["footer"].replace("<b>", "").replace("</b>", "")
                p_f.font.size = Pt(11)
                p_f.font.italic = True
                p_f.font.color.rgb = TEXT_MUTED

    prs.save(OUTPUT_PPTX)
    print(f"[OK] Master PPTX successfully generated: {OUTPUT_PPTX}")


# ===========================================================================
# PDF GENERATOR (REPORTLAB WIDESCREEN 16:9 - 960 x 540 PT)
# ===========================================================================
def build_pdf():
    """Generate the complete matching 30-page PDF deck."""
    PAGE_WIDTH = 960
    PAGE_HEIGHT = 540

    c = canvas.Canvas(str(OUTPUT_PDF), pagesize=(PAGE_WIDTH, PAGE_HEIGHT))

    for idx, data in enumerate(SLIDES_DATA):
        stype = data["type"]

        if stype == "COVER":
            c.setFillColor(RL_DARK_BG)
            c.rect(0, 0, PAGE_WIDTH, PAGE_HEIGHT, fill=1, stroke=0)

            c.setFillColor(RL_CYAN)
            c.roundRect(50, 440, 240, 24, 4, fill=1, stroke=0)
            c.setFillColor(RL_DARK_TEXT)
            c.setFont("Helvetica-Bold", 10)
            c.drawString(60, 447, data["tag"])

            c.setFont("Helvetica-Bold", 32)
            c.drawString(50, 385, data["title"])

            c.setFont("Helvetica", 14)
            c.setFillColor(RL_DARK_MUTED)
            c.drawString(50, 345, data["subtitle"])

            c.setFillColor(RL_DARK_CARD)
            c.roundRect(50, 180, 860, 120, 8, fill=1, stroke=0)

            c.setFillColor(RL_CYAN)
            c.setFont("Helvetica-Bold", 12)
            c.drawString(75, 260, data["meta"])

            c.setFillColor(RL_DARK_TEXT)
            c.setFont("Helvetica", 11)
            c.drawString(75, 230, data["author"])

            c.setFillColor(RL_GREEN)
            c.drawString(75, 205, "Open-Source Repository: https://github.com/miguelo0203/basketball-analytics")

        elif stype == "SECTION_DIVIDER":
            c.setFillColor(RL_DARK_BG)
            c.rect(0, 0, PAGE_WIDTH, PAGE_HEIGHT, fill=1, stroke=0)

            c.setFillColor(RL_CYAN)
            c.setFont("Helvetica-Bold", 64)
            c.drawString(60, 390, data["number"])

            c.setFillColor(RL_ORANGE)
            c.rect(60, 370, 140, 4, fill=1, stroke=0)

            c.setFillColor(RL_DARK_TEXT)
            c.setFont("Helvetica-Bold", 30)
            c.drawString(60, 310, data["title"])

            c.setFillColor(RL_DARK_MUTED)
            c.setFont("Helvetica", 14)
            c.drawString(60, 265, data["subtitle"])

        elif stype == "CLOSING":
            c.setFillColor(RL_DARK_BG)
            c.rect(0, 0, PAGE_WIDTH, PAGE_HEIGHT, fill=1, stroke=0)

            c.setFillColor(RL_GREEN)
            c.roundRect(50, 440, 240, 24, 4, fill=1, stroke=0)
            c.setFillColor(RL_DARK_TEXT)
            c.setFont("Helvetica-Bold", 10)
            c.drawString(60, 447, data["tag"])

            c.setFont("Helvetica-Bold", 30)
            c.drawString(50, 385, data["title"])

            c.setFont("Helvetica", 14)
            c.setFillColor(RL_DARK_MUTED)
            c.drawString(50, 345, data["subtitle"])

            c.setFillColor(RL_DARK_CARD)
            c.roundRect(50, 160, 860, 140, 8, fill=1, stroke=0)

            c.setFillColor(RL_CYAN)
            c.setFont("Helvetica-Bold", 13)
            c.drawString(75, 260, "GitHub Repository: " + data["github_url"])

            c.setFillColor(RL_GREEN)
            c.setFont("Helvetica-Bold", 12)
            c.drawString(75, 230, data["release"])

            c.setFillColor(RL_DARK_TEXT)
            c.setFont("Helvetica", 11)
            c.drawString(75, 200, data["citation"])

        else:
            # Light Content Background
            c.setFillColor(RL_LIGHT_BG)
            c.rect(0, 0, PAGE_WIDTH, PAGE_HEIGHT, fill=1, stroke=0)

            c.setFillColor(RL_CYAN)
            c.setFont("Helvetica-Bold", 9)
            c.drawString(50, 500, data.get("category", "INTERNATIONAL BASKETBALL ANALYTICS").upper())

            c.setFillColor(RL_TEXT_MAIN)
            c.setFont("Helvetica-Bold", 20)
            c.drawString(50, 470, data["title"])

            c.setFillColor(RL_TEXT_MUTED)
            c.setFont("Helvetica", 11)
            c.drawString(50, 448, data.get("subtitle", ""))

            # Layout Dispatcher
            if stype in ["EXEC_SUMMARY", "TAKEAWAYS"]:
                cards = data.get("cards", data.get("points", []))
                coords = [
                    (50, 260, 415, 165),
                    (495, 260, 415, 165),
                    (50, 75, 415, 165),
                    (495, 75, 415, 165)
                ]
                rl_colors = [RL_CYAN, RL_ORANGE, RL_GREEN, RL_PURPLE]
                for (title, desc, _), (x, y, w, h), col in zip(cards, coords, rl_colors):
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(x, y, w, h, 6, fill=1, stroke=1)

                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 12)
                    c.drawString(x + 15, y + h - 25, title)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica", 9.5)
                    words = desc.split()
                    lines = []
                    curr = []
                    for word in words:
                        if len(" ".join(curr + [word])) > 55:
                            lines.append(" ".join(curr))
                            curr = [word]
                        else:
                            curr.append(word)
                    if curr:
                        lines.append(" ".join(curr))

                    for j, line in enumerate(lines):
                        c.drawString(x + 15, y + h - 50 - j * 15, line)

            elif stype == "BIG_METRICS":
                metrics = data["metrics"]
                for i, (val, label, sub) in enumerate(metrics):
                    x = 50 + i * 220
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(x, 75, 205, 350, 6, fill=1, stroke=1)

                    c.setFillColor(RL_CYAN)
                    c.setFont("Helvetica-Bold", 28)
                    c.drawString(x + 15, 370, val)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica-Bold", 10)
                    c.drawString(x + 15, 340, label)

                    c.setFillColor(RL_TEXT_MUTED)
                    c.setFont("Helvetica", 9)
                    words = sub.split()
                    lines = []
                    curr = []
                    for word in words:
                        if len(" ".join(curr + [word])) > 25:
                            lines.append(" ".join(curr))
                            curr = [word]
                        else:
                            curr.append(word)
                    if curr:
                        lines.append(" ".join(curr))
                    for j, line in enumerate(lines):
                        c.drawString(x + 15, 310 - j * 14, line)

            elif stype in ["SPLIT_PROBLEM", "CASE_STUDY", "LIMITATIONS_MATRIX"]:
                # Left Column
                c.setFillColor(RL_CARD_BG)
                c.setStrokeColor(RL_CARD_BORDER)
                c.roundRect(50, 75, 420, 350, 6, fill=1, stroke=1)

                c.setFillColor(RL_ORANGE)
                c.setFont("Helvetica-Bold", 13)
                c.drawString(65, 395, data.get("left_title", "Analysis"))

                c.setFillColor(RL_TEXT_MAIN)
                c.setFont("Helvetica", 9.5)
                y_pos = 365
                for item in data.get("left_items", []):
                    words = item.split()
                    lines = []
                    curr = []
                    for word in words:
                        if len(" ".join(curr + [word])) > 52:
                            lines.append(" ".join(curr))
                            curr = [word]
                        else:
                            curr.append(word)
                    if curr:
                        lines.append(" ".join(curr))
                    for j, line in enumerate(lines):
                        prefix = "• " if j == 0 else "   "
                        c.drawString(65, y_pos, prefix + line)
                        y_pos -= 14
                    y_pos -= 6

                # Right Column
                c.setFillColor(RL_CARD_BG)
                c.setStrokeColor(RL_CARD_BORDER)
                c.roundRect(490, 75, 420, 350, 6, fill=1, stroke=1)

                c.setFillColor(RL_CYAN)
                c.setFont("Helvetica-Bold", 13)
                c.drawString(505, 395, data.get("right_title", "Solution"))

                c.setFillColor(RL_TEXT_MAIN)
                c.setFont("Helvetica", 9.5)
                y_pos = 365
                for item in data.get("right_items", []):
                    words = item.split()
                    lines = []
                    curr = []
                    for word in words:
                        if len(" ".join(curr + [word])) > 52:
                            lines.append(" ".join(curr))
                            curr = [word]
                        else:
                            curr.append(word)
                    if curr:
                        lines.append(" ".join(curr))
                    for j, line in enumerate(lines):
                        prefix = "• " if j == 0 else "   "
                        c.drawString(505, y_pos, prefix + line)
                        y_pos -= 14
                    y_pos -= 6

            elif stype == "FLOW_ARCHITECTURE":
                steps = data["steps"]
                rl_cols = [RL_CYAN, RL_ORANGE, RL_GREEN, RL_PURPLE]
                for i, (layer, name, desc, _) in enumerate(steps):
                    col = rl_cols[i % len(rl_cols)]
                    x = 50 + i * 220
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(x, 75, 205, 350, 6, fill=1, stroke=1)

                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 10)
                    c.drawString(x + 15, 395, layer)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica-Bold", 13)
                    c.drawString(x + 15, 370, name)

                    c.setFillColor(RL_TEXT_MUTED)
                    c.setFont("Helvetica", 9.5)
                    words = desc.split()
                    lines = []
                    curr = []
                    for word in words:
                        if len(" ".join(curr + [word])) > 26:
                            lines.append(" ".join(curr))
                            curr = [word]
                        else:
                            curr.append(word)
                    if curr:
                        lines.append(" ".join(curr))
                    for j, line in enumerate(lines):
                        c.drawString(x + 15, 335 - j * 15, line)

            elif stype == "ARCHETYPES_GRID":
                archetypes = data["archetypes"]
                coords = [
                    (50, 260, 270, 165),
                    (345, 260, 270, 165),
                    (640, 260, 270, 165),
                    (50, 75, 270, 165),
                    (345, 75, 270, 165),
                    (640, 75, 270, 165)
                ]
                rl_colors = [RL_CYAN, RL_ORANGE, RL_GREEN, RL_PURPLE, RL_CYAN, RL_ORANGE]
                for (name, desc, _), (x, y, w, h), col in zip(archetypes, coords, rl_colors):
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(x, y, w, h, 6, fill=1, stroke=1)

                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 12)
                    c.drawString(x + 12, y + h - 25, name)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica", 9)
                    words = desc.split()
                    lines = []
                    curr = []
                    for word in words:
                        if len(" ".join(curr + [word])) > 35:
                            lines.append(" ".join(curr))
                            curr = [word]
                        else:
                            curr.append(word)
                    if curr:
                        lines.append(" ".join(curr))
                    for j, line in enumerate(lines):
                        c.drawString(x + 12, y + h - 50 - j * 14, line)

            elif stype == "FOUR_FACTORS":
                factors = data["factors"]
                coords = [
                    (50, 260, 415, 165),
                    (495, 260, 415, 165),
                    (50, 75, 415, 165),
                    (495, 75, 415, 165)
                ]
                rl_colors = [RL_CYAN, RL_ORANGE, RL_GREEN, RL_PURPLE]
                for (name, subtitle, desc, formula, _), (x, y, w, h), col in zip(factors, coords, rl_colors):
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(x, y, w, h, 6, fill=1, stroke=1)

                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 12)
                    c.drawString(x + 15, y + h - 25, name + " (" + subtitle + ")")

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica", 9.5)
                    c.drawString(x + 15, y + h - 50, desc)

                    c.setFillColor(RL_CYAN)
                    c.setFont("Helvetica-Bold", 9)
                    c.drawString(x + 15, y + 20, formula)

            elif stype == "ML_CALIBRATION":
                metrics = data["metrics"]
                for i, (val, label, sub) in enumerate(metrics):
                    x = 50 + i * 295
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(x, 75, 275, 350, 6, fill=1, stroke=1)

                    c.setFillColor(RL_CYAN)
                    c.setFont("Helvetica-Bold", 32)
                    c.drawString(x + 20, 360, val)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(x + 20, 325, label)

                    c.setFillColor(RL_TEXT_MUTED)
                    c.setFont("Helvetica", 9.5)
                    words = sub.split()
                    lines = []
                    curr = []
                    for word in words:
                        if len(" ".join(curr + [word])) > 36:
                            lines.append(" ".join(curr))
                            curr = [word]
                        else:
                            curr.append(word)
                    if curr:
                        lines.append(" ".join(curr))
                    for j, line in enumerate(lines):
                        c.drawString(x + 20, 290 - j * 15, line)

            elif stype == "WALK_FORWARD":
                timeline = data["timeline"]
                for i, (fold, train, test) in enumerate(timeline):
                    y = 340 - i * 85
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(50, y, 860, 70, 6, fill=1, stroke=1)

                    c.setFillColor(RL_CYAN)
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(70, y + 42, fold)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica", 10)
                    c.drawString(140, y + 42, train)

                    c.setFillColor(RL_GREEN)
                    c.setFont("Helvetica-Bold", 10)
                    c.drawString(140, y + 20, "-> " + test)

            # MVP-34 POLISHED SLIDE 10: SHRINKAGE_VISUAL
            elif stype == "SHRINKAGE_VISUAL":
                # Left Column: 4-stage flow
                c.setFillColor(RL_CARD_BG)
                c.setStrokeColor(RL_CARD_BORDER)
                c.roundRect(50, 75, 415, 350, 6, fill=1, stroke=1)

                c.setFillColor(RL_CYAN)
                c.setFont("Helvetica-Bold", 12)
                c.drawString(65, 395, "ESTIMATION FLOW: FROM NOISE TO STABILITY")

                y_pos = 360
                rl_cols = [RL_ORANGE, RL_ORANGE, RL_CYAN, RL_GREEN]
                for i, (step_t, step_d, _) in enumerate(data["flow_steps"]):
                    col = rl_cols[i % len(rl_cols)]
                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 10.5)
                    prefix = "v  " if i > 0 else "*  "
                    c.drawString(65, y_pos, prefix + step_t)

                    c.setFillColor(RL_TEXT_MUTED)
                    c.setFont("Helvetica", 9)
                    c.drawString(80, y_pos - 16, step_d)
                    y_pos -= 65

                # Right Column: Top Callout Card
                c.setFillColor(RL_DARK_CARD)
                c.setStrokeColor(RL_CYAN)
                c.roundRect(495, 290, 415, 135, 6, fill=1, stroke=1)

                c_val, c_lbl, c_desc = data["callout"]
                c.setFillColor(RL_CYAN)
                c.setFont("Helvetica-Bold", 14)
                c.drawString(515, 390, c_val + "  |  " + c_lbl)

                c.setFillColor(RL_DARK_TEXT)
                c.setFont("Helvetica", 10)
                c.drawString(515, 360, "Bayesian shrinkage weight contracting noisy short-sample")
                c.drawString(515, 342, "shooting estimates toward 3,767 player career baselines.")
                c.setFillColor(RL_GREEN)
                c.setFont("Helvetica-Bold", 9.5)
                c.drawString(515, 310, "[OK] High-volume samples shrink less; short streaks regress.")

                # Middle Spectrum Card
                c.setFillColor(RL_CARD_BG)
                c.setStrokeColor(RL_CARD_BORDER)
                c.roundRect(495, 140, 415, 135, 6, fill=1, stroke=1)

                c.setFillColor(RL_TEXT_MAIN)
                c.setFont("Helvetica-Bold", 11)
                c.drawString(515, 248, "CONCEPTUAL SHRINKAGE SPECTRUM")

                for j, (label, desc, col) in enumerate(data["spectrum_items"]):
                    y_item = 218 - j * 32
                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 9.5)
                    c.drawString(515, y_item, "- " + label)
                    c.setFillColor(RL_TEXT_MUTED)
                    c.setFont("Helvetica", 9)
                    c.drawString(720, y_item, "-> " + desc)

                # Bottom Bootstrap Badge
                c.setFillColor(RL_CARD_BG)
                c.setStrokeColor(RL_CARD_BORDER)
                c.roundRect(495, 75, 415, 52, 6, fill=1, stroke=1)
                c.setFillColor(RL_GREEN)
                c.setFont("Helvetica-Bold", 9)
                c.drawString(510, 103, "Cluster Bootstrap (B = 5,000):")
                c.setFillColor(RL_TEXT_MAIN)
                c.setFont("Helvetica", 8.5)
                c.drawString(510, 88, "Non-parametric 95% confidence intervals across 3,767 qualified campaigns.")

            # MVP-34 POLISHED SLIDE 14: MONTE_CARLO_VISUAL
            elif stype == "MONTE_CARLO_VISUAL":
                # Top Pipeline Bar
                c.setFillColor(RL_DARK_CARD)
                c.setStrokeColor(RL_CYAN)
                c.roundRect(50, 385, 860, 40, 6, fill=1, stroke=1)
                c.setFillColor(RL_CYAN)
                c.setFont("Helvetica-Bold", 10.5)
                c.drawString(65, 402, "SIMULATION PIPELINE: " + data["pipeline"])

                # Left Panel: Bracket
                c.setFillColor(RL_CARD_BG)
                c.setStrokeColor(RL_CARD_BORDER)
                c.roundRect(50, 75, 415, 295, 6, fill=1, stroke=1)

                c.setFillColor(RL_TEXT_MAIN)
                c.setFont("Helvetica-Bold", 12)
                c.drawString(65, 345, "DYNAMIC TOURNAMENT BRACKET FLOW")

                y_b = 310
                for r_name, nodes in data["bracket_rounds"]:
                    c.setFillColor(RL_CYAN)
                    c.setFont("Helvetica-Bold", 10)
                    c.drawString(65, y_b, ">  " + r_name)
                    c.setFillColor(RL_TEXT_MUTED)
                    c.setFont("Helvetica", 9)
                    c.drawString(85, y_b - 16, "  |  ".join(nodes))
                    y_b -= 55

                # Right Panel: 3 Metric Cards
                rl_cols = [RL_CYAN, RL_ORANGE, RL_GREEN]
                for i, (val, lbl, desc, _) in enumerate(data["metrics"]):
                    col = rl_cols[i % len(rl_cols)]
                    y_m = 275 - i * 100
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(495, y_m, 415, 90, 6, fill=1, stroke=1)

                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 13)
                    c.drawString(515, y_m + 62, val + "  |  " + lbl)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica", 9)
                    words = desc.split()
                    lines = []
                    curr = []
                    for word in words:
                        if len(" ".join(curr + [word])) > 55:
                            lines.append(" ".join(curr))
                            curr = [word]
                        else:
                            curr.append(word)
                    if curr:
                        lines.append(" ".join(curr))
                    for j, line in enumerate(lines):
                        c.drawString(515, y_m + 38 - j * 14, line)

            # MVP-34 POLISHED SLIDE 15: WORKSPACE_MOCKUP
            elif stype == "WORKSPACE_MOCKUP":
                # App Header
                c.setFillColor(RL_DARK_CARD)
                c.setStrokeColor(RL_CYAN)
                c.roundRect(50, 385, 860, 40, 6, fill=1, stroke=1)
                c.setFillColor(RL_CYAN)
                c.setFont("Helvetica-Bold", 10.5)
                c.drawString(65, 402, "[APP]  " + data["app_title"])

                # Sidebar
                c.setFillColor(RL_DARK_BG)
                c.setStrokeColor(RL_DARK_CARD)
                c.roundRect(50, 75, 230, 295, 6, fill=1, stroke=1)

                c.setFillColor(RL_CYAN)
                c.setFont("Helvetica-Bold", 10)
                c.drawString(65, 345, "APP FILTERS")

                for i, item in enumerate(data["sidebar_items"]):
                    c.setFillColor(RL_DARK_TEXT)
                    c.setFont("Helvetica", 8.5)
                    c.drawString(65, 315 - i * 22, "- " + item)

                # Main Evidence: Available (Left) vs Locked (Right)
                c.setFillColor(RL_CARD_BG)
                c.setStrokeColor(RL_GREEN)
                c.roundRect(295, 145, 320, 225, 6, fill=1, stroke=1)

                c.setFillColor(RL_GREEN)
                c.setFont("Helvetica-Bold", 10)
                c.drawString(310, 345, "PRE-GAME EVIDENCE (AVAILABLE)")

                for i, item in enumerate(data["pre_game_items"]):
                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica", 8.5)
                    c.drawString(310, 315 - i * 28, "- " + item)

                # Locked Box
                c.setFillColor(RL_CARD_BG)
                c.setStrokeColor(RL_ORANGE)
                c.roundRect(630, 145, 280, 225, 6, fill=1, stroke=1)

                c.setFillColor(RL_ORANGE)
                c.setFont("Helvetica-Bold", 10)
                c.drawString(645, 345, "POST-GAME (QUARANTINED)")

                for i, item in enumerate(data["post_game_locked"]):
                    c.setFillColor(RL_TEXT_MUTED)
                    c.setFont("Helvetica", 8.5)
                    c.drawString(645, 315 - i * 28, item)

                # Bottom Timeline Ribbon
                c.setFillColor(RL_DARK_CARD)
                c.setStrokeColor(RL_CYAN)
                c.roundRect(295, 75, 615, 55, 6, fill=1, stroke=1)

                c.setFillColor(RL_DARK_TEXT)
                c.setFont("Helvetica-Bold", 9)
                timeline_str = "  ->  ".join([name for name, _ in data["timeline_steps"]])
                c.drawString(310, 98, "ANTI-HINDSIGHT AUDIT TRAIL:  " + timeline_str)

            # MVP-34 POLISHED SLIDE 27: ROLE_VALUE_TREE
            elif stype == "ROLE_VALUE_TREE":
                # Top Tree Header
                c.setFillColor(RL_DARK_CARD)
                c.setStrokeColor(RL_CYAN)
                c.roundRect(50, 385, 860, 40, 6, fill=1, stroke=1)
                c.setFillColor(RL_CYAN)
                c.setFont("Helvetica-Bold", 10.5)
                c.drawString(65, 402, "ORGANIZATIONAL IMPACT:  " + data["club_header"])

                # 3 Role Columns
                roles = data["roles"]
                rl_cols = [RL_CYAN, RL_ORANGE, RL_GREEN]
                for i, (role_name, role_q, items, _) in enumerate(roles):
                    col = rl_cols[i % len(rl_cols)]
                    x = 50 + i * 295
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(x, 140, 275, 230, 6, fill=1, stroke=1)

                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(x + 15, 345, role_name)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica-Bold", 10)
                    c.drawString(x + 15, 322, f'"{role_q}"')

                    for j, item in enumerate(items):
                        c.setFillColor(RL_TEXT_MUTED)
                        c.setFont("Helvetica", 8.5)
                        c.drawString(x + 15, 290 - j * 24, "- " + item)

                # Bottom Value Ribbon
                c.setFillColor(RL_DARK_CARD)
                c.setStrokeColor(RL_GREEN)
                c.roundRect(50, 75, 860, 50, 6, fill=1, stroke=1)
                c.setFillColor(RL_DARK_TEXT)
                c.setFont("Helvetica-Bold", 9.5)
                c.drawString(65, 96, "VALUE PIPELINE:  " + data["value_flow"])

            # MVP-34 POLISHED SLIDE 28: ROADMAP_TIMELINE
            elif stype == "ROADMAP_TIMELINE":
                # Top Timeline Header
                c.setFillColor(RL_DARK_CARD)
                c.setStrokeColor(RL_CYAN)
                c.roundRect(50, 385, 860, 40, 6, fill=1, stroke=1)
                c.setFillColor(RL_CYAN)
                c.setFont("Helvetica-Bold", 10.5)
                c.drawString(65, 402, "TRANSITION TIMELINE:  " + data["timeline_header"])

                # 3 Phase Cards
                phases = data["phases"]
                rl_cols = [RL_CYAN, RL_ORANGE, RL_GREEN]
                for i, (p_name, p_days, p_items, _) in enumerate(phases):
                    col = rl_cols[i % len(rl_cols)]
                    x = 50 + i * 295
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(x, 140, 275, 230, 6, fill=1, stroke=1)

                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(x + 15, 345, p_name)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica-Bold", 10)
                    c.drawString(x + 15, 322, p_days)

                    for j, item in enumerate(p_items):
                        c.setFillColor(RL_TEXT_MUTED)
                        c.setFont("Helvetica", 8.5)
                        words = item.split()
                        lines = []
                        curr = []
                        for word in words:
                            if len(" ".join(curr + [word])) > 36:
                                lines.append(" ".join(curr))
                                curr = [word]
                            else:
                                curr.append(word)
                        if curr:
                            lines.append(" ".join(curr))
                        for k, line in enumerate(lines):
                            prefix = "- " if k == 0 else "   "
                            c.drawString(x + 15, 290 - j * 32 - k * 12, prefix + line)

                # Bottom Milestone Card
                c.setFillColor(RL_DARK_CARD)
                c.setStrokeColor(RL_GREEN)
                c.roundRect(50, 75, 860, 50, 6, fill=1, stroke=1)
                c.setFillColor(RL_GREEN)
                c.setFont("Helvetica-Bold", 9.5)
                c.drawString(65, 96, "[*]  " + data["milestone"])

            elif stype == "CROSS_LANGUAGE":
                items = data["items"]
                rl_colors = [RL_CYAN, RL_GREEN, RL_ORANGE, RL_PURPLE]
                for i, (title, desc, _) in enumerate(items):
                    col = rl_colors[i % len(rl_colors)]
                    y = 340 - i * 85
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(50, y, 860, 70, 6, fill=1, stroke=1)

                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 11)
                    c.drawString(70, y + 42, title)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica", 10)
                    c.drawString(220, y + 42, desc)

            elif stype == "REPRODUCIBILITY":
                c.setFillColor(RL_DARK_CARD)
                c.setStrokeColor(RL_CYAN)
                c.roundRect(50, 360, 860, 55, 6, fill=1, stroke=1)

                c.setFillColor(RL_CYAN)
                c.setFont("Helvetica-Bold", 14)
                c.drawString(70, 382, "$ " + data["code"])

                steps = data["steps"]
                rl_cols = [RL_CYAN, RL_ORANGE, RL_GREEN, RL_PURPLE]
                for i, (step_name, step_desc, _) in enumerate(steps):
                    col = rl_cols[i % len(rl_cols)]
                    x = 50 + i * 220
                    c.setFillColor(RL_CARD_BG)
                    c.setStrokeColor(RL_CARD_BORDER)
                    c.roundRect(x, 75, 205, 265, 6, fill=1, stroke=1)

                    c.setFillColor(col)
                    c.setFont("Helvetica-Bold", 10)
                    c.drawString(x + 12, 310, step_name)

                    c.setFillColor(RL_TEXT_MAIN)
                    c.setFont("Helvetica", 9)
                    words = step_desc.split()
                    lines = []
                    curr = []
                    for word in words:
                        if len(" ".join(curr + [word])) > 26:
                            lines.append(" ".join(curr))
                            curr = [word]
                        else:
                            curr.append(word)
                    if curr:
                        lines.append(" ".join(curr))
                    for j, line in enumerate(lines):
                        c.drawString(x + 12, 280 - j * 14, line)

            else:
                cards = data.get("cards", data.get("roles", data.get("phases", [])))
                if cards:
                    rl_colors = [RL_CYAN, RL_ORANGE, RL_GREEN]
                    for i, item in enumerate(cards):
                        x = 50 + i * 295
                        c.setFillColor(RL_CARD_BG)
                        c.setStrokeColor(RL_CARD_BORDER)
                        c.roundRect(x, 75, 275, 350, 6, fill=1, stroke=1)

                        title = item[0]
                        desc = item[1]
                        col = rl_colors[i % len(rl_colors)]

                        c.setFillColor(col)
                        c.setFont("Helvetica-Bold", 12)
                        c.drawString(x + 18, 385, title)

                        c.setFillColor(RL_TEXT_MAIN)
                        c.setFont("Helvetica", 9.5)
                        words = desc.split()
                        lines = []
                        curr = []
                        for word in words:
                            if len(" ".join(curr + [word])) > 36:
                                lines.append(" ".join(curr))
                                curr = [word]
                            else:
                                curr.append(word)
                        if curr:
                            lines.append(" ".join(curr))
                        for j, line in enumerate(lines):
                            c.drawString(x + 18, 350 - j * 16, line)

            # Footer
            if "footer" in data:
                c.setFillColor(RL_TEXT_MUTED)
                c.setFont("Helvetica-Oblique", 9)
                clean_footer = data["footer"].replace("<b>", "").replace("</b>", "")
                c.drawString(50, 45, clean_footer)

            # Page numbering
            c.setFillColor(RL_TEXT_MUTED)
            c.setFont("Helvetica", 8)
            c.drawRightString(910, 45, f"{idx + 1} / {len(SLIDES_DATA)}")

        c.showPage()

    c.save()
    print(f"[OK] Master PDF successfully generated: {OUTPUT_PDF}")


if __name__ == "__main__":
    build_pptx()
    build_pdf()
