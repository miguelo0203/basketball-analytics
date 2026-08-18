"""Master Presentation Generator for International Basketball Analytics (2005-2024).

Generates the complete 30-slide professional PowerPoint presentation (.pptx)
focused on Analyst Workflow, Decision Support, and Methodological Rigor.
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

from src.config import PROJECT_ROOT


def build_complete_presentation():
    """Build the complete 30-slide presentation."""
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Color Palette Tokens
    NAVY_DARK = RGBColor(15, 23, 42)      # #0F172A
    SLATE_GRAY = RGBColor(71, 85, 105)    # #475569
    TEAL_ACCENT = RGBColor(2, 132, 199)   # #0284C7
    BG_CARD = RGBColor(248, 250, 252)     # #F8FAFC
    BORDER_CARD = RGBColor(226, 232, 240) # #E2E8F0
    WHITE = RGBColor(255, 255, 255)
    GREEN_ACCENT = RGBColor(22, 163, 74)  # #16A34A
    AMBER_ACCENT = RGBColor(217, 119, 6)  # #D97706
    RED_ACCENT = RGBColor(220, 38, 38)    # #DC2626
    BG_TEAL_LIGHT = RGBColor(240, 249, 255)

    def add_blank_slide():
        blank_layout = prs.slide_layouts[6]
        return prs.slides.add_slide(blank_layout)

    def add_header(slide, title_text, category_text="INTERNATIONAL BASKETBALL ANALYTICS (2005–2024)"):
        """Add standardized header banner to slide."""
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.35))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_top = tf_cat.margin_right = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(11)
        p_cat.font.bold = True
        p_cat.font.color.rgb = TEAL_ACCENT

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.72), Inches(11.7), Inches(0.8))
        tf = title_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.size = Pt(21)
        p.font.bold = True
        p.font.color.rgb = NAVY_DARK

    def add_card(slide, left, top, width, height, title=None, bg_color=BG_CARD, border_color=BORDER_CARD):
        """Add visual card container."""
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)

        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.4))
            tf = tb.text_frame
            tf.word_wrap = True
            tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(13)
            p.font.bold = True
            p.font.color.rgb = NAVY_DARK
        return shape

    # =========================================================================
    # SLIDE 1: PORTADA
    # =========================================================================
    s1 = add_blank_slide()
    bg = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY_DARK
    bg.line.color.rgb = NAVY_DARK

    tb1 = s1.shapes.add_textbox(Inches(1.0), Inches(1.8), Inches(11.333), Inches(4.0))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p1 = tf1.paragraphs[0]
    p1.text = "INTERNATIONAL BASKETBALL ANALYTICS"
    p1.font.size = Pt(36)
    p1.font.bold = True
    p1.font.color.rgb = WHITE
    
    p2 = tf1.add_paragraph()
    p2.text = "De los datos a la evidencia para apoyar decisiones de baloncesto"
    p2.font.size = Pt(20)
    p2.font.color.rgb = TEAL_ACCENT
    p2.space_before = Pt(14)

    p3 = tf1.add_paragraph()
    p3.text = "2005–2024  |  Basketball Data Analyst Portfolio  |  Miguel"
    p3.font.size = Pt(14)
    p3.font.color.rgb = RGBColor(148, 163, 184)
    p3.space_before = Pt(40)

    # =========================================================================
    # SLIDE 2: LA IDEA EN 30 SEGUNDOS
    # =========================================================================
    s2 = add_blank_slide()
    add_header(s2, "La Idea en 30 Segundos: Convertir Complejidad en Evidencia")
    
    add_card(s2, Inches(0.8), Inches(1.8), Inches(11.733), Inches(2.2), "El Propósito Fundamental")
    tb2 = s2.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(11.1), Inches(1.5))
    tf2 = tb2.text_frame
    tf2.word_wrap = True
    p2_1 = tf2.paragraphs[0]
    p2_1.text = "Construí un sistema analítico reproducible que combina datos de partido, rendimiento individual, contexto competitivo, validación táctica en vídeo y modelado estadístico para transformar información heterogénea en evidencia accionable que un entrenador o director deportivo pueda interpretar y utilizar con confianza."
    p2_1.font.size = Pt(15)
    p2_1.font.color.rgb = NAVY_DARK

    steps = [
        ("1. DATOS", "Actas oficiales FIBA y hashes SHA-256"),
        ("2. ANÁLISIS", "Four Factors, ritmo y roles funcionales"),
        ("3. VALIDACIÓN", "17 Folds walk-forward y vídeo (κ=0.80)"),
        ("4. CONTEXTO", "Aislamiento temporal anti-hindsight"),
        ("5. DECISIÓN", "Briefs prepartido con preguntas clave")
    ]
    for i, (title, desc) in enumerate(steps):
        left_pos = Inches(0.8 + i * 2.4)
        add_card(s2, left_pos, Inches(4.3), Inches(2.2), Inches(2.4), title)
        tb_step = s2.shapes.add_textbox(left_pos + Inches(0.15), Inches(4.8), Inches(1.9), Inches(1.7))
        tf_step = tb_step.text_frame
        tf_step.word_wrap = True
        p_st = tf_step.paragraphs[0]
        p_st.text = desc
        p_st.font.size = Pt(12)
        p_st.font.color.rgb = SLATE_GRAY

    # =========================================================================
    # SLIDE 3: QUÉ PROBLEMA RESUELVE UN ANALISTA
    # =========================================================================
    s3 = add_blank_slide()
    add_header(s3, "¿Qué Problema Resuelve un Analista de Baloncesto?")

    add_card(s3, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "El Desafío del Cuerpo Técnico")
    tb3_1 = s3.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf3_1 = tb3_1.text_frame
    tf3_1.word_wrap = True
    bullets_in = [
        "Sobrecarga de boxscores y números descontextualizados.",
        "Muestras pequeñas de torneo (6–9 partidos) con alta varianza de tiro.",
        "Sesgo retrospectivo (evaluar la decisión por el resultado final).",
        "Brecha entre analítica cuantitativa y la pizarra de los entrenadores.",
        "Falta de cuantificación honesta de la incertidumbre."
    ]
    for b in bullets_in:
        p = tf3_1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    add_card(s3, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "El Valor Real del Analista", bg_color=BG_TEAL_LIGHT, border_color=TEAL_ACCENT)
    tb3_2 = s3.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf3_2 = tb3_2.text_frame
    tf3_2.word_wrap = True
    bullets_out = [
        "Filtrar el ruido estadístico mediante baselines robustos.",
        "Unir las métricas avanzadas (Four Factors) con la cinta de vídeo.",
        "Superficiar contradicciones ocultas entre datos y sensaciones.",
        "Entregar briefs prepartido concisos (1.5 páginas, 2.5 min de lectura).",
        "Formular preguntas tácticas concretas para el cuerpo técnico."
    ]
    for b in bullets_out:
        p = tf3_2.add_paragraph()
        p.text = "✔ " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 4: QUÉ CONSTRUÍ (ARQUITECTURA)
    # =========================================================================
    s4 = add_blank_slide()
    add_header(s4, "Qué Construí: Arquitectura Integral Desacoplada en 9 Capas")

    layers = [
        ("1. Data Ingestion", "Parseo de actas FIBA y firmas SHA-256"),
        ("2. Data Quality & DB", "DuckDB relacional con 12 tablas normalizadas"),
        ("3. Basketball Analytics", "Four Factors de Oliver, Pace y Net Rating"),
        ("4. Player Analytics", "K-Means++ y PCA en 3.767 campañas (K=6)"),
        ("5. Tactical Validation", "Vídeo doblemente codificado (420 clips, κ=0.80)"),
        ("6. Machine Learning", "LightGBM con 17 folds walk-forward (ECE=0.0314)"),
        ("7. Tournament Simulation", "180.000 iteraciones Monte Carlo con shrinkage"),
        ("8. Decision Synthesis", "Matriz de 8 capas y motor de contradicciones"),
        ("9. Operational Workspace", "Streamlit interactivo con aislamiento anti-hindsight")
    ]
    for i, (name, detail) in enumerate(layers):
        row = i // 3
        col = i % 3
        l = Inches(0.8 + col * 4.0)
        t = Inches(1.8 + row * 1.7)
        add_card(s4, l, t, Inches(3.7), Inches(1.5), name)
        tb = s4.shapes.add_textbox(l + Inches(0.15), t + Inches(0.55), Inches(3.4), Inches(0.8))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = detail
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE_GRAY

    # =========================================================================
    # SLIDE 5: ESCALA DEL PROYECTO
    # =========================================================================
    s5 = add_blank_slide()
    add_header(s5, "Escala y Cobertura Histórica Verificada (2005–2024)")

    metrics_s5 = [
        ("18", "Torneos Internacionales", "EuroBasket, Copa del Mundo FIBA y JJ.OO."),
        ("1,145", "Partidos Oficiales", "Cobertura longitudinal completa de 20 años"),
        ("27,353", "Actuaciones de Jugador", "Registros individuales en fact_player_game"),
        ("3,767", "Campañas Cualificadas", "Jugadores con >= 40 minutos en el torneo"),
        ("6", "Arquetipos Funcionales", "Roles minados por K-Means++ y PCA"),
        ("17 Folds", "Walk-Forward Folds", "1.105 partidos evaluados out-of-sample"),
        ("180k", "Simulaciones Monte Carlo", "10.000 iteraciones por campeonato"),
        ("195", "Tests Automatizados", "100% de tasa de éxito en suite pytest")
    ]
    for i, (val, title, sub) in enumerate(metrics_s5):
        row = i // 4
        col = i % 4
        l = Inches(0.8 + col * 3.0)
        t = Inches(1.8 + row * 2.5)
        add_card(s5, l, t, Inches(2.8), Inches(2.2))
        tb = s5.shapes.add_textbox(l + Inches(0.15), t + Inches(0.15), Inches(2.5), Inches(1.9))
        tf = tb.text_frame
        tf.word_wrap = True
        
        p1 = tf.paragraphs[0]
        p1.text = val
        p1.font.size = Pt(26)
        p1.font.bold = True
        p1.font.color.rgb = TEAL_ACCENT

        p2 = tf.add_paragraph()
        p2.text = title
        p2.font.size = Pt(13)
        p2.font.bold = True
        p2.font.color.rgb = NAVY_DARK
        p2.space_before = Pt(4)

        p3 = tf.add_paragraph()
        p3.text = sub
        p3.font.size = Pt(10)
        p3.font.color.rgb = SLATE_GRAY
        p3.space_before = Pt(4)

    # =========================================================================
    # SLIDE 6: CALIDAD DE DATOS Y PIPELINE DE QA
    # =========================================================================
    s6 = add_blank_slide()
    add_header(s6, "Calidad de Datos: 'Un Análisis Sofisticado con Malos Datos es Incorrecto'")

    add_card(s6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "Pipeline de QA Determinista")
    tb6_1 = s6.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf6_1 = tb6_1.text_frame
    tf6_1.word_wrap = True
    steps_qa = [
        "1. Ingesta Bruta & Hash SHA-256: Congelación inmutable del origen.",
        "2. Resolución de Entidades: 2.124 jugadores unificados sin duplicados tipográficos.",
        "3. Regla de Minutos (200 min): Cuadre exacto de tiempo de juego.",
        "4. Regla de Marcador: Suma de puntos individual = tanteo final del partido.",
        "5. Coherencia de Posesiones: Bilateralidad de posesiones entre equipos."
    ]
    for s in steps_qa:
        p = tf6_1.add_paragraph()
        p.text = s
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(12)

    add_card(s6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "Esquema Relacional DuckDB OLAP", bg_color=BG_TEAL_LIGHT)
    tb6_2 = s6.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf6_2 = tb6_2.text_frame
    tf6_2.word_wrap = True
    p_duck = tf6_2.paragraphs[0]
    p_duck.text = "Estructura de Almacén Normalizado (12 Tablas):"
    p_duck.font.size = Pt(13)
    p_duck.font.bold = True
    p_duck.font.color.rgb = NAVY_DARK
    
    tables_list = [
        "• dim_tournament, dim_team, dim_player, dim_coach",
        "• fact_game (1.145 partidos, fechas, marcadores, ritmo)",
        "• fact_team_game (2.290 observaciones avanzadas)",
        "• fact_player_game (27.353 registros individuales)",
        "• mart_player_roles (3.767 campañas clusterizadas)",
        "• mart_tactical_video (420 posesiones codificadas)"
    ]
    for t in tables_list:
        p = tf6_2.add_paragraph()
        p.text = t
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE_GRAY
        p.space_after = Pt(8)

    # =========================================================================
    # SLIDE 7: POR QUÉ NO USÉ SOLO PPG
    # =========================================================================
    s7 = add_blank_slide()
    add_header(s7, "Más Allá de los Puntos por Partido (PPG): El Valor Funcional")

    add_card(s7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "La Trampa del Volumen Simple (PPG)")
    tb7_1 = s7.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf7_1 = tb7_1.text_frame
    tf7_1.word_wrap = True
    bullets_ppg = [
        "Un jugador puede anotar 18 PPG con bajo True Shooting (TS% < 48%) y alto uso de posesiones, restando eficiencia al colectivo.",
        "Los PPG acumulados no consideran el ritmo de posesiones del equipo.",
        "Ignoran la generación indirecta (tiros creados, bloqueos, espaciado).",
        "Ocultan el impacto defensivo y las pérdidas provocadas."
    ]
    for b in bullets_ppg:
        p = tf7_1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(12)

    add_card(s7, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "La Descomposición Analítica Completa", bg_color=BG_TEAL_LIGHT)
    tb7_2 = s7.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf7_2 = tb7_2.text_frame
    tf7_2.word_wrap = True
    bullets_full = [
        "✔ Eficiencia de Tiro: eFG% y True Shooting% ponderando el valor del triple.",
        "✔ Ajuste por Posesiones: Estadísticas normalizadas por 40 minutos.",
        "✔ Cuatro Factores de Oliver: Impacto en rebote, pérdidas y faltas.",
        "✔ Gravedad Espacial: Atracción de marcas y liberaciones en pick-and-pop.",
        "✔ Net Rating On/Off: Rendimiento del quinteto con y sin el jugador."
    ]
    for b in bullets_full:
        p = tf7_2.add_paragraph()
        p.text = b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(12)

    # =========================================================================
    # SLIDE 8: FOUR FACTORS Y CONTEXTO
    # =========================================================================
    s8 = add_blank_slide()
    add_header(s8, "Los Four Factors de Dean Oliver como Preguntas Tácticas")

    factors = [
        ("Effective Field Goal % (eFG%)", "¿Estamos perdiendo eficiencia por mala selección de tiro o por falta de finalización interior?"),
        ("Turnover Rate (TOV%)", "¿Nuestras pérdidas proceden de pases arriesgados en transición o de colapso ante presión en media pista?"),
        ("Offensive Rebound % (ORB%)", "¿Generamos segundas opciones sin comprometer el balance defensivo de repliegue?"),
        ("Free Throw Rate (FTR)", "¿Estamos atacando el aro y cargando de faltas al pívot rival en drop o abusando del tiro exterior?")
    ]
    for i, (name, q) in enumerate(factors):
        row = i // 2
        col = i % 2
        l = Inches(0.8 + col * 6.0)
        t = Inches(1.8 + row * 2.5)
        add_card(s8, l, t, Inches(5.6), Inches(2.2), name)
        tb = s8.shapes.add_textbox(l + Inches(0.2), t + Inches(0.65), Inches(5.2), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"Pregunta táctica para el cuerpo técnico:\n\"{q}\""
        p.font.size = Pt(12)
        p.font.italic = True
        p.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 9: PLAYER ANALYTICS & ESTABILIDAD
    # =========================================================================
    s9 = add_blank_slide()
    add_header(s9, "Player Analytics: No Cuánto Produce, sino Qué Función Desempeña")

    add_card(s9, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "Normalización y Contexto de Uso")
    tb9_1 = s9.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf9_1 = tb9_1.text_frame
    tf9_1.word_wrap = True
    bullets_pa = [
        "Estadísticas Per-40 Minutos: Comparación equitativa neutralizando rotaciones cortas.",
        "Usage Rate (USG%): Volumen de posesiones finalizadas por tiro, falta o pérdida.",
        "Ratio Asistencias / Pérdidas: Control de balón del generador primario.",
        "Estabilidad Longitudinal: Resistencia del perfil individual a través de múltiples torneos."
    ]
    for b in bullets_pa:
        p = tf9_1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(12)

    add_card(s9, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "El Cambio de Paradigma", bg_color=BG_TEAL_LIGHT)
    tb9_2 = s9.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf9_2 = tb9_2.text_frame
    tf9_2.word_wrap = True
    p9_2 = tf9_2.paragraphs[0]
    p9_2.text = "El análisis cuantitativo moderno no responde a la pregunta de si un jugador 'es bueno' en términos absolutos."
    p9_2.font.size = Pt(13)
    p9_2.font.color.rgb = NAVY_DARK
    
    p9_3 = tf9_2.add_paragraph()
    p9_3.text = "Responde a:\n\n1. ¿Qué rol estadístico desempeña en pista?\n2. ¿Qué quintetos maximizan su eficiencia neta?\n3. ¿Cómo compensa la estructura del equipo sus puntos débiles?"
    p9_3.font.size = Pt(13)
    p9_3.font.bold = True
    p9_3.font.color.rgb = TEAL_ACCENT
    p9_3.space_before = Pt(14)

    # =========================================================================
    # SLIDE 10: FUNCTIONAL ARCHETYPES (K-MEANS++ & PCA)
    # =========================================================================
    s10 = add_blank_slide()
    add_header(s10, "Arquetipos Funcionales: Minería No Supervisada en 3.767 Campañas")

    archetypes = [
        ("1. Primary Initiator", "Alto uso de balón, generación tras bote y tiros en aclarado."),
        ("2. Movement Spacer", "Triples tras recepción, juego indirecto y gravedad perimetral."),
        ("3. Interior Hub", "Postes dominantes, rebote ofensivo y distribución interior."),
        ("4. Floor General", "Bases directores, control de tempo y ratio AST/TOV de élite."),
        ("5. Defensive Anchor", "Pívots intimidadores, protección de aro y bloqueo directo."),
        ("6. Balanced Wing", "Aleros complementarios con producción bidireccional equilibrada.")
    ]
    for i, (name, desc) in enumerate(archetypes):
        row = i // 3
        col = i % 3
        l = Inches(0.8 + col * 4.0)
        t = Inches(1.8 + row * 2.5)
        add_card(s10, l, t, Inches(3.7), Inches(2.2), name)
        tb = s10.shapes.add_textbox(l + Inches(0.15), t + Inches(0.65), Inches(3.4), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE_GRAY

    # =========================================================================
    # SLIDE 11: DEL NÚMERO A LA PISTA (PUENTE CUANTI-CUALI)
    # =========================================================================
    s11 = add_blank_slide()
    add_header(s11, "Del Número a la Pista: El Puente entre Estadística y Vídeo")

    bridge_steps = [
        ("1. Estadística", "Descubre un patrón agregado en los Four Factors (ej. alta eficiencia en tiros de media distancia)."),
        ("2. Hipótesis", "Plantea qué ajuste táctico rival está generando esa anomalía numérica."),
        ("3. Evidencia Táctica", "Audita la cinta de vídeo bajo protocolo estandarizado (coberturas de bloqueo directo)."),
        ("4. Interpretación", "Traduce el hallazgo a una pregunta accionable para la sesión de vídeo del entrenador.")
    ]
    for i, (st_name, st_desc) in enumerate(bridge_steps):
        l = Inches(0.8 + i * 3.0)
        add_card(s11, l, Inches(1.8), Inches(2.8), Inches(4.9), st_name)
        tb = s11.shapes.add_textbox(l + Inches(0.15), Inches(2.5), Inches(2.5), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = st_desc
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 12: VALIDACIÓN TÁCTICA EN VÍDEO (κ = 0.80)
    # =========================================================================
    s12 = add_blank_slide()
    add_header(s12, "Validación Táctica: Protocolo y Fiabilidad Inter-Evaluador")

    add_card(s12, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "Protocolo de Doble Codificación")
    tb12_1 = s12.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf12_1 = tb12_1.text_frame
    tf12_1.word_wrap = True
    bullets_v = [
        "Muestra: 420 posesiones observadas en 36 partidos clave.",
        "Variables: Profundidad de drop coverage en P&R y contestación en closeouts.",
        "Doble Codificación: Dos analistas independientes evaluando las mismas acciones.",
        "Fiabilidad Tipo de Acción: Cohen's Kappa κ = 1.00 (acuerdo total).",
        "Fiabilidad Calificación Defensiva: Cohen's Kappa κ = 0.80 (acuerdo sustancial)."
    ]
    for b in bullets_v:
        p = tf12_1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    add_card(s12, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "El Valor de la Fiabilidad Científica", bg_color=BG_TEAL_LIGHT)
    tb12_2 = s12.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf12_2 = tb12_2.text_frame
    tf12_2.word_wrap = True
    p12_2 = tf12_2.paragraphs[0]
    p12_2.text = "\"Si una observación táctica depende exclusivamente de quién mira el vídeo, es imposible convertirla en evidencia reproducible para un cuerpo técnico.\""
    p12_2.font.size = Pt(14)
    p12_2.font.italic = True
    p12_2.font.color.rgb = NAVY_DARK
    
    p12_3 = tf12_2.add_paragraph()
    p12_3.text = "La capa cualitativa del sistema no es un sustituto del ojo del entrenador: es un protocolo estandarizado para generar hipótesis tácticas contrastables."
    p12_3.font.size = Pt(12)
    p12_3.font.color.rgb = SLATE_GRAY
    p12_3.space_before = Pt(16)

    # =========================================================================
    # SLIDE 13: MACHINE LEARNING SUPERVISADO
    # =========================================================================
    s13 = add_blank_slide()
    add_header(s13, "Machine Learning Supervisado: Benchmark y Rendimiento Histórico")

    models_bench = [
        ("Naive Baseline", "Brier: 0.2450", "AUC: 0.5000", "MAE: 14.82 pts"),
        ("Logistic Regression", "Brier: 0.2085", "AUC: 0.7320", "MAE: 12.45 pts"),
        ("ElasticNet Regressor", "Brier: 0.2040", "AUC: 0.7410", "MAE: 12.10 pts"),
        ("LightGBM Calibrado", "Brier: 0.1967", "AUC: 0.7613", "MAE: 11.739 pts")
    ]
    for i, (m_name, brier, auc, mae) in enumerate(models_bench):
        l = Inches(0.8 + i * 3.0)
        bg_col = BG_TEAL_LIGHT if "LightGBM" in m_name else BG_CARD
        border_col = TEAL_ACCENT if "LightGBM" in m_name else BORDER_CARD
        add_card(s13, l, Inches(1.8), Inches(2.8), Inches(3.2), m_name, bg_color=bg_col, border_color=border_col)
        tb = s13.shapes.add_textbox(l + Inches(0.15), Inches(2.5), Inches(2.5), Inches(2.3))
        tf = tb.text_frame
        tf.word_wrap = True
        for stat in [brier, auc, mae]:
            p = tf.add_paragraph()
            p.text = stat
            p.font.size = Pt(13)
            p.font.bold = True if "LightGBM" in m_name else False
            p.font.color.rgb = NAVY_DARK
            p.space_after = Pt(6)

    # Bottom caveat card
    add_card(s13, Inches(0.8), Inches(5.3), Inches(11.733), Inches(1.4), "Disciplina de Comunicación")
    tb13_bot = s13.shapes.add_textbox(Inches(1.0), Inches(5.75), Inches(11.3), Inches(0.8))
    tf13_bot = tb13_bot.text_frame
    tf13_bot.word_wrap = True
    p_c = tf13_bot.paragraphs[0]
    p_c.text = "Estos números describen el rendimiento predictivo retrospectivo bajo este protocolo temporal. No constituyen una afirmación de que el modelo 'adivine' partidos con certeza absoluta."
    p_c.font.size = Pt(12)
    p_c.font.color.rgb = SLATE_GRAY

    # =========================================================================
    # SLIDE 14: POR QUÉ EL WALK-FORWARD IMPORTA
    # =========================================================================
    s14 = add_blank_slide()
    add_header(s14, "Validación Walk-Forward: Cero Fuga de Información del Futuro")

    add_card(s14, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "17 Folds Temporales Expansivos")
    tb14_1 = s14.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf14_1 = tb14_1.text_frame
    tf14_1.word_wrap = True
    folds_text = [
        "Fold 1: Entrena 2005 (T1) ➔ Evalúa 2006 (T2)",
        "Fold 2: Entrena 2005–2006 (T1..T2) ➔ Evalúa 2007 (T3)",
        "Fold 3: Entrena 2005–2007 (T1..T3) ➔ Evalúa 2008 (T4)",
        "...",
        "Fold 17: Entrena 2005–2023 (T1..T17) ➔ Evalúa 2024 (T18)",
        "Total Evaluado Out-of-Sample: 1.105 partidos."
    ]
    for f in folds_text:
        p = tf14_1.add_paragraph()
        p.text = f
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(8)

    add_card(s14, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "Principio de Honestidad Metodológica", bg_color=BG_TEAL_LIGHT)
    tb14_2 = s14.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf14_2 = tb14_2.text_frame
    tf14_2.word_wrap = True
    p14_2 = tf14_2.paragraphs[0]
    p14_2.text = "\"El modelo jamás ve el futuro para construir sus predicciones prepartido.\""
    p14_2.font.size = Pt(15)
    p14_2.font.bold = True
    p14_2.font.color.rgb = TEAL_ACCENT
    
    p14_3 = tf14_2.add_paragraph()
    p14_3.text = "\nEn analítica profesional, prefiero una métrica temporal imperfecta pero honesta a una precisión inflada artificialmente mediante data leakage o mezclas aleatorias (k-fold shuffle)."
    p14_3.font.size = Pt(13)
    p14_3.font.color.rgb = NAVY_DARK
    p14_3.space_before = Pt(10)

    # =========================================================================
    # SLIDE 15: INCERTIDUMBRE Y CALIBRACIÓN
    # =========================================================================
    s15 = add_blank_slide()
    add_header(s15, "Incertidumbre y Calibración: Cuánto Confiar en lo que Dice el Modelo")

    add_card(s15, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "Calibración Isotónica (ECE = 0.0314)")
    tb15_1 = s15.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf15_1 = tb15_1.text_frame
    tf15_1.word_wrap = True
    bullets_cal = [
        "Expected Calibration Error (ECE): 0.0314 (3.14% de error medio en probabilidad).",
        "Brier Score Out-of-Sample: 0.1967 frente al 0.2500 de un modelo aleatorio.",
        "Significado Práctico: Cuando el sistema asigna un 70% de victoria, el equipo gana empíricamente 7 de cada 10 veces.",
        "Seguridad de Simulación: Las probabilidades son seguras para propagar a Monte Carlo."
    ]
    for b in bullets_cal:
        p = tf15_1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    add_card(s15, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "Inferencia y Control de Varianza", bg_color=BG_TEAL_LIGHT)
    tb15_2 = s15.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf15_2 = tb15_2.text_frame
    tf15_2.word_wrap = True
    bullets_inf = [
        "Bootstrap por Conglomerados (B=5.000): Intervalos de confianza al 95% preservando correlaciones de torneo.",
        "Permutation Testing (P=10.000): Contrastes de hipótesis no paramétricos sin supuestos gaussianos.",
        "Corrección FDR (Benjamini-Hochberg): Control de tasa de falsos descubrimientos en comparaciones múltiples."
    ]
    for b in bullets_inf:
        p = tf15_2.add_paragraph()
        p.text = "✔ " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(12)

    # =========================================================================
    # SLIDE 16: QUÉ SIGNIFICA SHAP (ATRIBUCIÓN != CAUSALIDAD)
    # =========================================================================
    s16 = add_blank_slide()
    add_header(s16, "Atribución TreeSHAP: 'Estas Asociaciones son Predictivas, No Causales'")

    add_card(s16, Inches(0.8), Inches(1.8), Inches(11.733), Inches(1.4), "Aviso Fundamental de Epistemología Analítica", bg_color=RGBColor(254, 242, 242), border_color=RED_ACCENT)
    tb16_top = s16.shapes.add_textbox(Inches(1.0), Inches(2.25), Inches(11.3), Inches(0.8))
    tf16_top = tb16_top.text_frame
    tf16_top.word_wrap = True
    p16_t = tf16_top.paragraphs[0]
    p16_t.text = "ESTAS ASOCIACIONES SON PREDICTIVAS, NO CAUSALES. Los valores SHAP explican cómo el algoritmo pondera estadísticamente las variables históricas; no garantizan que una intervención en pista produzca ese efecto."
    p16_t.font.size = Pt(13)
    p16_t.font.bold = True
    p16_t.font.color.rgb = RED_ACCENT

    shap_vars = [
        ("Net Rating Diferencial", "Mayor contribuyente global al desplazamiento log-odds de probabilidad."),
        ("Diferencial de eFG%", "Impacto masivo de la eficiencia de tiro efectiva en torneos cortos."),
        ("Forma en el Torneo", "Ponderación del margen de puntos en la fase de grupos reciente."),
        ("Turnover Rate (TOV%)", "Castigo severo por pérdidas vivas que generan contraataques.")
    ]
    for i, (vname, vdesc) in enumerate(shap_vars):
        l = Inches(0.8 + i * 3.0)
        add_card(s16, l, Inches(3.5), Inches(2.8), Inches(3.2), vname)
        tb = s16.shapes.add_textbox(l + Inches(0.15), Inches(4.2), Inches(2.5), Inches(2.3))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = vdesc
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE_GRAY

    # =========================================================================
    # SLIDE 17: MONTE CARLO
    # =========================================================================
    s17 = add_blank_slide()
    add_header(s17, "Simulaciones Monte Carlo: 180.000 Iteraciones de Torneo")

    add_card(s17, Inches(0.8), Inches(1.8), Inches(11.733), Inches(2.2), "El Concepto de la Simulación")
    tb17 = s17.shapes.add_textbox(Inches(1.1), Inches(2.3), Inches(11.1), Inches(1.5))
    tf17 = tb17.text_frame
    tf17.word_wrap = True
    p17_1 = tf17.paragraphs[0]
    p17_1.text = "PROBABILIDADES CALIBRADAS ➔ 10.000 REPLAYS POR CAMPEONATO ➔ DISTRIBUCIÓN DE AVANCE\n\nNo intentamos afirmar que 'España tiene exactamente un X% fijo de ganar el oro'. Afirmamos que bajo estas probabilidades calibradas por partido y este formato de eliminatorias, esta es la distribución empírica de escenarios simulados."
    p17_1.font.size = Pt(14)
    p17_1.font.color.rgb = NAVY_DARK

    sim_features = [
        ("Shrinkage de Probabilidad (λ=0.75)", "Ajuste conservador para evitar sobreconfianza en favoritos teóricos."),
        ("Propagación de Cuadros", "Cálculo de emparejamientos dinámicos según cruces de grupos."),
        ("Captura de Oro Top-1: 72.2%", "Consistencia retrospectiva en 13 de 18 torneos históricos."),
        ("Captura de Medalla Top-4: 100%", "El medallista real siempre se situó en el Top-4 proyectado.")
    ]
    for i, (title, desc) in enumerate(sim_features):
        l = Inches(0.8 + i * 3.0)
        add_card(s17, l, Inches(4.3), Inches(2.8), Inches(2.4), title)
        tb = s17.shapes.add_textbox(l + Inches(0.15), Inches(5.0), Inches(2.5), Inches(1.5))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = Pt(11)
        p.font.color.rgb = SLATE_GRAY

    # =========================================================================
    # SLIDE 18: POR QUÉ SIMULAR UN TORNEO
    # =========================================================================
    s18 = add_blank_slide()
    add_header(s18, "¿Por Qué Simular un Torneo? Partido Individual vs. Campeonato")

    add_card(s18, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "El Partido Individual")
    tb18_1 = s18.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf18_1 = tb18_1.text_frame
    tf18_1.word_wrap = True
    bullets_p = [
        "Un partido es un suceso de alta varianza puntual.",
        "Un mal día en el tiro exterior (20% en triples) puede costar la eliminación.",
        "El foco táctico es 100% inmediato: emparejamientos y coberturas específicas."
    ]
    for b in bullets_p:
        p = tf18_1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(14)

    add_card(s18, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "El Torneo Acumulado (Simulación)", bg_color=BG_TEAL_LIGHT)
    tb18_2 = s18.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf18_2 = tb18_2.text_frame
    tf18_2.word_wrap = True
    bullets_t = [
        "Una victoria o derrota en fase de grupos altera el lado del cuadro.",
        "Permite evaluar la probabilidad acumulada de alcanzar Cuartos, Semifinales y Final.",
        "Ayuda a directores deportivos a planificar profundidad de plantilla para torneos de 8 partidos en 15 días."
    ]
    for b in bullets_t:
        p = tf18_2.add_paragraph()
        p.text = "✔ " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(14)

    # =========================================================================
    # SLIDE 19: DEL ANÁLISIS A LA DECISIÓN (6 CAPAS)
    # =========================================================================
    s19 = add_blank_slide()
    add_header(s19, "Del Análisis a la Decisión: 6 Capas de Evidencia Integrada")

    six_layers = [
        ("1. Producción / Eficiencia", "Four Factors, True Shooting y ritmo."),
        ("2. Fiabilidad Estadística", "Muestras de minutos e intervalos bootstrap."),
        ("3. Rol Funcional", "Arquetipos y equilibrio posicional en pista."),
        ("4. Evidencia Táctica", "Drop coverage en P&R y velocidad de closeouts."),
        ("5. Impacto Predictivo", "Probabilidad calibrada y margen esperado ML."),
        ("6. Contexto de Torneo", "Simulación de cuadro y opciones de medalla.")
    ]
    for i, (lname, ldesc) in enumerate(six_layers):
        row = i // 3
        col = i % 3
        l = Inches(0.8 + col * 4.0)
        t = Inches(1.8 + row * 2.5)
        add_card(s19, l, t, Inches(3.7), Inches(2.2), lname)
        tb = s19.shapes.add_textbox(l + Inches(0.15), t + Inches(0.65), Inches(3.4), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ldesc
        p.font.size = Pt(13)
        p.font.color.rgb = SLATE_GRAY

    # =========================================================================
    # SLIDE 20: CONTRADICCIONES (ESTADÍSTICA VS VÍDEO)
    # =========================================================================
    s20 = add_blank_slide()
    add_header(s20, "Detección de Contradicciones: 'No Oculto el Conflicto, lo Hago Explícito'")

    add_card(s20, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.9), "Ejemplo Real de Contradicción Táctica")
    tb20 = s20.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.0))
    tf20 = tb20.text_frame
    tf20.word_wrap = True

    p20_1 = tf20.paragraphs[0]
    p20_1.text = "BOXSCORE / NÚMEROS AGREGADOS: POSITIVO\nEE. UU. promediaba +31.2 puntos de margen en el torneo y el modelo otorgaba 73.2% de favoritismo."
    p20_1.font.size = Pt(13)
    p20_1.font.color.rgb = NAVY_DARK

    p20_2 = tf20.add_paragraph()
    p20_2.text = "CINTA DE VÍDEO TÁCTICO: NEGATIVO / VULNERABLE\nEl pívot defensor de EE. UU. (Howard) jugaba drop profundo por debajo del tiro libre, concediendo tiros liberados en pick-and-pop a pívots exteriores (Pau y Marc Gasol, Garbajosa)."
    p20_2.font.size = Pt(13)
    p20_2.font.color.rgb = RED_ACCENT
    p20_2.space_before = Pt(12)

    p20_3 = tf20.add_paragraph()
    p20_3.text = "ACCIÓN DEL ANALISTA:\nSuperficiar la discrepancia en el Brief Prepartido: Formular la pregunta de si los bases españoles pueden castigar el drop mediante pick-and-pop sin forzar pases interiores arriesgados."
    p20_3.font.size = Pt(13)
    p20_3.font.bold = True
    p20_3.font.color.rgb = TEAL_ACCENT
    p20_3.space_before = Pt(14)

    # =========================================================================
    # SLIDE 21: EJEMPLO REAL: PEKÍN 2008
    # =========================================================================
    s21 = add_blank_slide()
    add_header(s21, "Caso Real: Final Olímpica de Pekín 2008 (España vs. EE. UU.)")

    add_card(s21, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "El Estado Prepartido (T-1)")
    tb21_1 = s21.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf21_1 = tb21_1.text_frame
    tf21_1.word_wrap = True
    bullets_21_1 = [
        "Precedente: Derrota por 37 puntos (119–82) en fase de grupos.",
        "Probabilidad Modelo: P(ESP) = 26.8% (Margen esperado: -8.5 pts).",
        "Hallazgo Four Factors: España tenía Net Rating +4.2 en media pista estática.",
        "Alerta Táctica: Explotar el pick-and-pop y usar zona 2-3 para frenar la transición."
    ]
    for b in bullets_21_1:
        p = tf21_1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    add_card(s21, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "Resultado y Pensamiento Probabilístico", bg_color=BG_TEAL_LIGHT)
    tb21_2 = s21.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf21_2 = tb21_2.text_frame
    tf21_2.word_wrap = True
    bullets_21_2 = [
        "Marcador Real: USA 118 – ESP 107 (11 puntos de diferencia).",
        "Desarrollo: España aplicó la zona 2-3 y el pick-and-pop, situándose a 4 pts a falta de 2:20.",
        "Incertidumbre: El margen de 11 pts cayó dentro del intervalo bootstrap 95% ([-16.8, +1.2]).",
        "Conclusión: El resultado real no invalida la probabilidad previa: es una realización posible dentro de la distribución."
    ]
    for b in bullets_21_2:
        p = tf21_2.add_paragraph()
        p.text = "✔ " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 22: EJEMPLO REAL: EUROBASKET 2015
    # =========================================================================
    s22 = add_blank_slide()
    add_header(s22, "Caso Real: EuroBasket 2015 (No Sobrerreaccionar a Muestras Cortas)")

    add_card(s22, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "La Crisis de Fase de Grupos")
    tb22_1 = s22.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf22_1 = tb22_1.text_frame
    tf22_1.word_wrap = True
    bullets_22_1 = [
        "Contexto: España sufrió dos derrotas en grupos (vs Serbia e Italia).",
        "Alarma Mediática: Percepción de fin de ciclo y crisis competitiva.",
        "Realidad Numérica: Las derrotas se produjeron por varianza extrema de tiro rival (Italia encestó 57% en triples) manteniendo España un Net Rating positivo."
    ]
    for b in bullets_22_1:
        p = tf22_1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(12)

    add_card(s22, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "La Proyección del Sistema", bg_color=BG_TEAL_LIGHT)
    tb22_2 = s22.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf22_2 = tb22_2.text_frame
    tf22_2.word_wrap = True
    bullets_22_2 = [
        "Estabilidad del Modelo: El sistema mantuvo a España como favorito a medalla (67.6% título implícito).",
        "Eje Central: Dominancia absoluta de Pau Gasol como Interior Hub.",
        "Desenlace Real: España ganó el Oro venciendo a Francia y Lituania.",
        "Lección: El analista aporta valor evitando que el staff sobrerreaccione a muestras pequeñas de tiro."
    ]
    for b in bullets_22_2:
        p = tf22_2.add_paragraph()
        p.text = "✔ " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(12)

    # =========================================================================
    # SLIDE 23: QUÉ PUEDE HACER UN ENTRENADOR CON ESTO
    # =========================================================================
    s23 = add_blank_slide()
    add_header(s23, "¿Qué Puede Hacer un Entrenador con Esta Información?")

    coach_qs = [
        ("1. Detección de Fugas", "¿Dónde estamos perdiendo eficiencia neta por posesión?"),
        ("2. Control de Pérdidas", "¿Qué tipo de posesiones generan nuestras pérdidas y permiten contraataques?"),
        ("3. Sostenibilidad de Quintetos", "¿Qué perfiles de jugador están sosteniendo la producción real?"),
        ("4. Ajustes Tácticos", "¿Qué cambios defensivos son consistentes con el drop del rival?"),
        ("5. Sensibilidad de Cuadro", "¿Qué escenarios del torneo son especialmente sensibles a un cruce?"),
        ("6. Calidad de Evidencia", "¿Qué conclusiones tienen respaldo sólido y cuáles son ruido estadístico?")
    ]
    for i, (qtitle, qtext) in enumerate(coach_qs):
        row = i // 3
        col = i % 3
        l = Inches(0.8 + col * 4.0)
        t = Inches(1.8 + row * 2.5)
        add_card(s23, l, t, Inches(3.7), Inches(2.2), qtitle)
        tb = s23.shapes.add_textbox(l + Inches(0.15), t + Inches(0.65), Inches(3.4), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = f"\"{qtext}\""
        p.font.size = Pt(13)
        p.font.italic = True
        p.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 24: EJEMPLO DE BRIEF PARA STAFF
    # =========================================================================
    s24 = add_blank_slide()
    add_header(s24, "Ejemplo Real de Brief Prepartido para el Cuerpo Técnico")

    add_card(s24, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.9), "Estructura Estandarizada de 1.5 Páginas (Lectura: 2.5 Minutos)")
    tb24 = s24.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.0))
    tf24 = tb24.text_frame
    tf24.word_wrap = True

    sections_brief = [
        ("1. CONTEXTO Y RITMO", "Ritmo estimado (72.1 poss), probabilidad base (26.8%) e intervalo empírico."),
        ("2. TRES PRIORIDADES", "Control de pérdidas en primera línea, forzar tiro exterior y cargar el rebote."),
        ("3. CONTRADICCIÓN TÁCTICA", "Discrepancia entre favoritismo numérico y vulnerabilidad en drop de P&R."),
        ("4. PREGUNTAS PARA LA PIZARRA", "¿Cómo castigamos el espacio del tráiler? ¿Aplicamos zona 2-3 tras canasta?"),
        ("5. LÍMITES DEL INFORME", "La estadística no anticipa rachas de tiro individuales de 8 metros.")
    ]
    for stitle, sdesc in sections_brief:
        p = tf24.add_paragraph()
        p.text = f"{stitle}: {sdesc}"
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 25: QUÉ HARÍA DIFERENTE EN UN CLUB REAL
    # =========================================================================
    s25 = add_blank_slide()
    add_header(s25, "¿Qué Haría Diferente en un Entorno de Club Profesional?")

    add_card(s25, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "Este Proyecto (Demostración Metodológica)")
    tb25_1 = s25.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf25_1 = tb25_1.text_frame
    tf25_1.word_wrap = True
    bullets_dem = [
        "Datos históricos de selecciones nacionales (2005–2024).",
        "Boxscores oficiales y codificación cualitativa de vídeo.",
        "Torneos cortos de 15 días (6–9 partidos).",
        "Demuestra metodología, ingeniería y pensamiento crítico."
    ]
    for b in bullets_dem:
        p = tf25_1.add_paragraph()
        p.text = "• " + b
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(12)

    add_card(s25, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "En un Club Real (Despliegue Operativo)", bg_color=BG_TEAL_LIGHT)
    tb25_2 = s25.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf25_2 = tb25_2.text_frame
    tf25_2.word_wrap = True
    bullets_real = [
        "✔ Feeds en vivo de tracking óptico 25Hz (Second Spectrum).",
        "✔ Integración con play-by-play detallado (Synergy Sports).",
        "✔ Datos de carga física y biometría (Catapult GPS / Firstbeat).",
        "✔ Base de datos de contratos y mercado para restricciones salariales.",
        "✔ Ritmo semanal de 2-3 partidos con briefs inmediatos post-shootaround."
    ]
    for b in bullets_real:
        p = tf25_2.add_paragraph()
        p.text = b
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 26: QUÉ APORTARÍA COMO ANALISTA
    # =========================================================================
    s26 = add_blank_slide()
    add_header(s26, "Qué Aportaría al Departamento de Baloncesto")

    contributions = [
        ("DATA ENGINEERING", "Construcción y mantenimiento de almacenes DuckDB/SQL inmutables y pipelines de QA."),
        ("ANALYTICS & ML", "Métricas avanzadas, calibración de modelos y simulaciones Monte Carlo robustas."),
        ("BASKETBALL DOMAIN", "Traducción fluida de números a conceptos de juego (P&R, Four Factors, espaciado)."),
        ("VALIDATION & QA", "Control estricto de fuga de datos, estimación de incertidumbre e integridad."),
        ("COMMUNICATION", "Informes prepartido concisos para personas que no necesitan ver código."),
        ("DECISION SUPPORT", "Transformación de evidencia en opciones tácticas y preguntas de pizarra.")
    ]
    for i, (cname, cdesc) in enumerate(contributions):
        row = i // 3
        col = i % 3
        l = Inches(0.8 + col * 4.0)
        t = Inches(1.8 + row * 2.5)
        add_card(s26, l, t, Inches(3.7), Inches(2.2), cname)
        tb = s26.shapes.add_textbox(l + Inches(0.15), t + Inches(0.65), Inches(3.4), Inches(1.4))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = cdesc
        p.font.size = Pt(12)
        p.font.color.rgb = SLATE_GRAY

    # =========================================================================
    # SLIDE 27: LO QUE NO HARÍA (HUMILDAD PROFESIONAL)
    # =========================================================================
    s27 = add_blank_slide()
    add_header(s27, "Lo Que NO Haría: Límites y Madurez Profesional")

    add_card(s27, Inches(0.8), Inches(1.8), Inches(11.733), Inches(4.9), "Líneas Rojas Profesionales")
    tb27 = s27.shapes.add_textbox(Inches(1.1), Inches(2.4), Inches(11.1), Inches(4.0))
    tf27 = tb27.text_frame
    tf27.word_wrap = True

    no_list = [
        "❌ NO sustituiría la autoridad ni el criterio del entrenador principal.",
        "❌ NO presentaría una correlación estadística como una relación causal garantizada.",
        "❌ NO vendería una probabilidad predictiva como una certeza matemática.",
        "❌ NO ocultaría la incertidumbre estadística ni los intervalos de varianza de tiro.",
        "❌ NO utilizaría datos del futuro para simular conocimiento retrospectivo (anti-hindsight).",
        "❌ NO asumiría que un modelo histórico funciona automáticamente en una liga diferente sin calibración local."
    ]
    for item in no_list:
        p = tf27.add_paragraph()
        p.text = item
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    # =========================================================================
    # SLIDE 28: LIMITACIONES Y HONESTIDAD
    # =========================================================================
    s28 = add_blank_slide()
    add_header(s28, "Límites Metodológicos y Alcance del Proyecto")

    limits = [
        ("Muestras Cortas", "6–9 partidos por torneo generan varianza natural en porcentajes de tiro."),
        ("Sin Tracking 25Hz", "El proyecto no dispone de coordenadas continuas XYZ de cámaras en vivo."),
        ("Rotación de Plantilla", "Las selecciones cambian de convocatoria cada verano a diferencia de clubes."),
        ("Diferencias de Eras", "Impacto de cambios de reglas FIBA (retraso del triple a 6.75m en 2010).")
    ]
    for i, (lname, ldesc) in enumerate(limits):
        l = Inches(0.8 + i * 3.0)
        add_card(s28, l, Inches(1.8), Inches(2.8), Inches(4.9), lname)
        tb = s28.shapes.add_textbox(l + Inches(0.15), Inches(2.5), Inches(2.5), Inches(4.0))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = ldesc
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK

    # =========================================================================
    # SLIDE 29: EL REPOSITORIO GITHUB
    # =========================================================================
    s29 = add_blank_slide()
    add_header(s29, "Estructura del Repositorio e Itinerarios de Revisión")

    add_card(s29, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.9), "Estructura Limpia del Repositorio")
    tb29_1 = s29.shapes.add_textbox(Inches(1.0), Inches(2.4), Inches(5.2), Inches(4.0))
    tf29_1 = tb29_1.text_frame
    tf29_1.word_wrap = True
    repo_dirs = [
        "README.md ➔ Puerta de entrada en español.",
        "docs/ ➔ Suite completa de 12 documentos técnicos.",
        "portfolio/ ➔ Caso flagship y guía de figuras.",
        "src/analytics/ ➔ 10 Módulos de analítica y Streamlit.",
        "data/ ➔ Almacén DuckDB y Parquet marts.",
        "tests/ ➔ 21 Módulos pytest (195 tests pasando)."
    ]
    for d in repo_dirs:
        p = tf29_1.add_paragraph()
        p.text = "📁 " + d
        p.font.size = Pt(12)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(10)

    add_card(s29, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.9), "Itinerarios de Lectura Optimizados", bg_color=BG_TEAL_LIGHT)
    tb29_2 = s29.shapes.add_textbox(Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.0))
    tf29_2 = tb29_2.text_frame
    tf29_2.word_wrap = True
    journeys = [
        "⏱️ 2 Minutos ➔ README.md + Resumen en 30s.",
        "⏱️ 5 Minutos ➔ Caso Flagship Pekín 2008 + Demo Streamlit.",
        "⏱️ 15 Minutos ➔ docs/machine_learning.md + Vídeo.",
        "⏱️ 30+ Minutos ➔ Inspección de código SQL/DuckDB y tests."
    ]
    for j in journeys:
        p = tf29_2.add_paragraph()
        p.text = j
        p.font.size = Pt(13)
        p.font.color.rgb = NAVY_DARK
        p.space_after = Pt(14)

    # =========================================================================
    # SLIDE 30: CIERRE Y FILOSOFÍA DEL ANALISTA
    # =========================================================================
    s30 = add_blank_slide()
    bg30 = s30.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg30.fill.solid()
    bg30.fill.fore_color.rgb = NAVY_DARK
    bg30.line.color.rgb = NAVY_DARK

    tb30 = s30.shapes.add_textbox(Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.8))
    tf30 = tb30.text_frame
    tf30.word_wrap = True

    p30_1 = tf30.paragraphs[0]
    p30_1.text = "MI OBJETIVO COMO ANALISTA"
    p30_1.font.size = Pt(32)
    p30_1.font.bold = True
    p30_1.font.color.rgb = WHITE

    p30_2 = tf30.add_paragraph()
    p30_2.text = "\"Convertir datos y evidencia de baloncesto en información clara, reproducible y útil para ayudar a tomar mejores decisiones.\""
    p30_2.font.size = Pt(20)
    p30_2.font.italic = True
    p30_2.font.color.rgb = TEAL_ACCENT
    p30_2.space_before = Pt(20)

    p30_3 = tf30.add_paragraph()
    p30_3.text = "DATOS  ➔  ANÁLISIS  ➔  EVIDENCIA  ➔  CONTEXTO  ➔  SOPORTE A DECISIONES"
    p30_3.font.size = Pt(15)
    p30_3.font.bold = True
    p30_3.font.color.rgb = RGBColor(148, 163, 184)
    p30_3.space_before = Pt(35)

    p30_4 = tf30.add_paragraph()
    p30_4.text = "El modelo no toma la decisión. Ayuda a que la decisión tenga mejor información detrás."
    p30_4.font.size = Pt(14)
    p30_4.font.color.rgb = WHITE
    p30_4.space_before = Pt(25)

    # Save presentation
    output_dir = PROJECT_ROOT / "portfolio" / "presentation"
    output_dir.mkdir(parents=True, exist_ok=True)
    pptx_path = output_dir / "International_Basketball_Analytics_Presentation.pptx"
    prs.save(pptx_path)
    print(f"Master Presentation generated successfully: {len(prs.slides)} slides -> {pptx_path}")
    return pptx_path


if __name__ == "__main__":
    build_complete_presentation()
