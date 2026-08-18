"""MVP-9 Automated PowerPoint Presentation Generator.

Generates reports/presentation/mvp9_analyst_portfolio.pptx (40 slides)
with dark sports-analytics theme, structured layouts, and embedded speaker notes.
"""

from pathlib import Path
import pptx
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

from src.config import REPORTS_DIR


class PresentationBuilder:
    """Automates production of the 40-slide master portfolio deck."""

    def __init__(self, output_path: Path = REPORTS_DIR / "presentation" / "mvp9_analyst_portfolio.pptx"):
        self.output_path = output_path
        self.output_path.parent.mkdir(parents=True, exist_ok=True)
        self.prs = pptx.Presentation()
        
        # 16:9 Widescreen dimensions
        self.prs.slide_width = Inches(13.333)
        self.prs.slide_height = Inches(7.5)

        # Palette definition
        self.COLOR_BG = RGBColor(15, 23, 42)        # Slate 900
        self.COLOR_CARD = RGBColor(30, 41, 59)      # Slate 800
        self.COLOR_ACCENT = RGBColor(16, 185, 129)  # Emerald 500
        self.COLOR_TEXT_MAIN = RGBColor(248, 250, 252) # Slate 50
        self.COLOR_TEXT_MUTED = RGBColor(148, 163, 184) # Slate 400
        self.COLOR_CYAN = RGBColor(56, 189, 248)    # Sky 400

    def add_slide(self, section: str, title: str, subtitle: str, bullets: list, notes: str = ""):
        """Add a formatted slide with structured cards and speaker notes."""
        blank_layout = self.prs.slide_layouts[6]
        slide = self.prs.slides.add_slide(blank_layout)

        # 1. Background Fill
        bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, self.prs.slide_width, self.prs.slide_height)
        bg.fill.solid()
        bg.fill.fore_color.rgb = self.COLOR_BG
        bg.line.fill.background()

        # 2. Top Header Container
        header_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.5), Inches(11.733), Inches(1.3))
        tf = header_box.text_frame
        tf.word_wrap = True
        tf.margin_left = tf.margin_top = tf.margin_right = tf.margin_bottom = 0

        # Section Tag
        p_sec = tf.paragraphs[0]
        p_sec.text = section.upper()
        p_sec.font.size = Pt(11)
        p_sec.font.bold = True
        p_sec.font.color.rgb = self.COLOR_ACCENT

        # Title
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = self.COLOR_TEXT_MAIN

        # Subtitle
        if subtitle:
            p_sub = tf.add_paragraph()
            p_sub.text = subtitle
            p_sub.font.size = Pt(13)
            p_sub.font.color.rgb = self.COLOR_CYAN

        # 3. Main Content Card
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.8), Inches(2.0), Inches(11.733), Inches(4.8))
        card.fill.solid()
        card.fill.fore_color.rgb = self.COLOR_CARD
        card.line.color.rgb = RGBColor(51, 65, 85)

        content_box = slide.shapes.add_textbox(Inches(1.2), Inches(2.3), Inches(10.933), Inches(4.2))
        c_tf = content_box.text_frame
        c_tf.word_wrap = True
        c_tf.margin_left = c_tf.margin_top = c_tf.margin_right = c_tf.margin_bottom = 0

        for i, b in enumerate(bullets):
            p = c_tf.paragraphs[0] if i == 0 else c_tf.add_paragraph()
            p.text = f"• {b}"
            p.font.size = Pt(15)
            p.font.color.rgb = self.COLOR_TEXT_MAIN
            p.space_after = Pt(12)

        # 4. Speaker Notes
        if notes:
            slide.notes_slide.notes_text_frame.text = notes

        return slide

    def build_full_presentation(self):
        """Construct all 40 slides across the 12 sections."""
        slides_data = [
            # SECTION 1
            ("Section 1: Why This Project Exists", "International Basketball Historical Analytics (2005–2025)", "Transforming 20 Years of Heterogeneous Basketball Data into Uncertainty-Aware Decision Support", [
                "Core Premise: The value of the analyst is not making the decision, but giving decision-makers better evidence and clearer uncertainty.",
                "Epistemological Foundation: Prediction != Causation | Observation != Proof | Model Output != Analyst Judgment.",
                "Comprehensive Scope: 18 Tournaments, 1,145 Games, 4,350 Campaigns, 128 Automated Unit Tests (100% Pass Rate).",
                "Presentation Objective: Demonstrating professional quantitative decision support for head coaches and sporting directors."
            ], "Welcome. Today I present a complete 20-year decision support stack for international basketball."),

            ("Section 1: Why This Project Exists", "The Problem of Uncertainty in Basketball Decisions", "Beyond Traditional Counting Stats and Outcome Bias", [
                "Small Sample Trap: International tournaments feature only 5 to 9 games per team over a 2-week span.",
                "Outcome Bias: A flawed tactical process can accidentally win on a buzzer-beater; a great process can suffer bad shooting variance.",
                "Traditional Boxscore Flaws: Points and rebounds ignore pace, possession efficiency, opponent quality, and defensive scheme context.",
                "Analyst Mission: Decouple process quality from random outcome noise through rigorous multi-layer modeling."
            ], "In tournament basketball, tiny samples create massive outcome bias. The analyst's job is evaluating underlying process quality."),

            ("Section 1: Why This Project Exists", "The Multi-Layer Evidence Hierarchy", "What Does a Coaching Staff and Sporting Leadership Actually Need?", [
                "Efficiency Context: True Shooting Efficiency (TS%), Turnover Rate (TOV%), and Four Factors possession impact.",
                "Statistical Reliability: Explicit Bootstrap Confidence Intervals (B=5,000) and sample exposure thresholds.",
                "Functional Roles: Categorizing players by tactical archetype rather than obsolete 1-to-5 position labels.",
                "Film Validation: Double-coded qualitative video analysis (Cohen's Kappa = 0.80) to inspect decision mechanics.",
                "Predictive & Simulation Support: Out-of-sample win probability shifts and Monte Carlo bracket advancement odds."
            ], "A coach cannot make roster decisions from an arbitrary single rating. We deliver evidence across 6 auditable layers."),

            ("Section 1: Why This Project Exists", "The Central Research Question", "Transforming Historical Basketball Evidence into Decision Support", [
                "Primary Question: 'How can an analyst transform historical match, player, and tactical data into defensible decision support?'",
                "Secondary Objectives: Decision Quality, Evidence Rigor, Probability Calibration, and Historical Generalization.",
                "Focus on Senior Men's FIBA Competitions: EuroBasket, FIBA World Cup, and Olympic Basketball Tournaments (2005–2024).",
                "End-to-End Execution: Full stack spanning Data Engineering, Econometrics, Supervised ML, Simulation, and Decision Science."
            ], "Our central question asks how heterogeneous data can be structured into robust decision support."),

            # SECTION 2
            ("Section 2: Project Scope & Universe", "Historical Scope & Tournament Coverage Timeline", "18 Certified Senior Men's Tournaments (2005–2024)", [
                "EuroBasket Universe (8 Tournaments): 2005, 2007, 2009, 2011, 2013, 2015, 2017, 2022.",
                "FIBA World Cup Universe (5 Tournaments): 2006, 2010, 2014, 2019, 2023.",
                "Olympic Basketball Universe (5 Tournaments): 2008, 2012, 2016, 2020, 2024.",
                "Temporal Fidelity: Reflects exact historical calendar cycles without fabricating unplayed tournament years."
            ], "Our dataset covers 18 premier international tournaments across 20 years with complete historical fidelity."),

            ("Section 2: Project Scope & Universe", "Verified Data Cardinalities & Warehouse Integrity", "Enterprise-Grade Relational Analytical Warehouse (DuckDB)", [
                "Canonical Match Universe: 1,145 Games (100% reconciled against official FIBA records with 0 missing games).",
                "Team-Game Observations: 2,290 Records across bilateral home/away perspectives.",
                "Player Campaigns: 4,350 Total Campaigns (3,767 qualified under the >= 40 minutes screening filter).",
                "Granular Player-Games: 27,353 Individual game boxscores and per-possession rate metrics.",
                "Tactical Film Layer: 420 Double-coded high-leverage possessions (Cohen's Kappa = 1.00 / 0.80)."
            ], "Every number in this presentation is backed by our audited relational database with zero missing records."),

            ("Section 2: Project Scope & Universe", "Functional Player Archetypes (Beyond Nominal Positions)", "Unsupervised Role Discovery via K-Means++ and PCA", [
                "Primary Initiator / Floor General: High usage, elite assist rate (AST% >= 25%), primary pick-and-roll ball-handler.",
                "Two-Way Scoring Wing / Slasher: Secondary creation, rim penetration, perimeter defensive versatility.",
                "Perimeter Movement Shooter / Spacer: High 3-point attempt rate (3PAr >= 50%), off-ball gravity and relocation.",
                "Stretch Big / Pick-and-Pop Forward: Frontcourt floor spacing, trailing transition threes, pick-and-pop execution.",
                "Low-Block Anchor / Interior Scorer: Post-up efficiency, offensive rebounding gravity, interior hub passing.",
                "Rim Protector / Roll Threat & Anchor: Defensive rim deterrence, pick-and-roll dive threat, screen setting."
            ], "Traditional position tags are obsolete. We cluster 4,350 campaigns into 6 functional archetypes."),

            # SECTION 3
            ("Section 3: Analytical Architecture", "The Full End-to-End Analytical Pipeline", "From Raw Immutable Data to Executive Decision Support", [
                "Stage 1-2: Data Acquisition & Warehouse Engineering (SHA-256 raw storage, DuckDB relational schema).",
                "Stage 3-4: Descriptive Analytics & Archetype Clustering (Four Factors, ITS Econometrics, K-Means++).",
                "Stage 5-6: Tactical Film & Supervised ML (Double-coded video IRR, 17-fold temporal walk-forward validation).",
                "Stage 7-8: Monte Carlo Simulation & Decision Dossiers (180k tournament runs, 6-layer candidate dossiers).",
                "Stage 9: Coaching Staff Decision Authority (Human expert judgment empowered by structured evidence)."
            ], "This architecture diagram maps our entire engineering and data science pipeline from raw data to coach handoff."),

            ("Section 3: Analytical Architecture", "Cumulative Stage Mapping (MVP-0 through MVP-8)", "How Each Analytical Layer Builds Directly on Upstream Foundations", [
                "MVP-0 to MVP-1: Guaranteed bitwise data engineering integrity and 100% tournament coverage.",
                "MVP-2: Quantified the macroeconomic impact of FIBA's 2010 3-point line expansion via Interrupted Time Series.",
                "MVP-3 to MVP-4: Ingested 27,353 player-games and engineered the multi-dimensional Candidate Fit Index (CFI).",
                "MVP-5 to MVP-6: Structured qualitative film coding and built calibrated pre-game predictive models (ECE = 0.0314).",
                "MVP-7 to MVP-8: Propagated probabilities through tournament brackets and generated auditable decision dossiers."
            ], "Each MVP stage was designed as a necessary prerequisite for the next, forming an integrated research stack."),

            ("Section 3: Analytical Architecture", "Methodological Guardrails & Zero Hindsight Invariant", "Eliminating Data Leakage and Retrospective Contamination", [
                "Strict Temporal Isolation: Pre-game features for game G strictly query matches occurring prior to game G's date.",
                "Bilateral Symmetry: Balanced bilateral representation prevents artificial home/away model bias.",
                "Retrospective Outcome Isolation: Tournament winners and medals are strictly isolated as evaluation benchmarks.",
                "No Data Overwrite: Raw provenance remains immutable with SHA-256 cryptographic verification."
            ], "We enforce strict temporal boundaries so models never have access to future information."),

            # SECTION 4
            ("Section 4: Tactical Film Evidence", "Quantitative Data vs Qualitative Film", "The Complementary Relationship in High-Performance Analytics", [
                "Structured Quantitative Data: Measures WHAT happened (shooting efficiency, turnover rate, possession margins).",
                "Qualitative Video Film: Investigates HOW and WHY it happened (decision speed, passing reads, defensive rotation angles).",
                "The Analytical Principle: Never make a major roster decision on boxscores alone; never make a decision on highlight clips alone.",
                "Integration: Video coding converts subjective scout notes into structured categorical variables."
            ], "Data tells us what happened; video tells us why. An elite analyst synthesizes both."),

            ("Section 4: Tactical Film Evidence", "Inter-Rater Reliability (IRR) & Video Coding Protocol", "Structuring Video as Scientifically Auditable Evidence", [
                "Double-Coded Protocol: 420 high-leverage possessions independently coded across multiple analysts.",
                "Action Types: Pick-and-Roll Reads, Closeout Attack/Defense, Post Hub Gravity, and Drop Coverage Navigation.",
                "Action Classification Concordance: Cohen's Kappa = 1.00 (100% perfect tactical agreement).",
                "Execution Quality Concordance: Cohen's Kappa = 0.80 (Substantial inter-rater reliability)."
            ], "We double-coded 420 possessions to prove that video observations can achieve scientific reliability."),

            ("Section 4: Tactical Film Evidence", "The Tactical Evidence Hierarchy & Contradiction Alerts", "Flagging Discrepancies Between Boxscores and Tape", [
                "Level 1: Direct Video Observation of tactical execution (double-coded, high IRR).",
                "Level 2: Empirical Rate Metrics (True Shooting, Turnover%, Rebound%).",
                "Level 3: Model-Derived Attributions (Predicted Net Impact, Centroid Distance).",
                "Contradiction Alert: Automatically triggered when boxscore stats contradict film execution quality."
            ], "Our evidence hierarchy automatically flags tactical contradictions when stats and film disagree."),

            # SECTION 5
            ("Section 5: Supervised Modeling & ML", "Temporal Walk-Forward Validation Framework", "Why Random K-Fold Splitting is Statistically Invalid in Sports", [
                "The Flaw of Random Splitting: Randomly shuffling tournament games leaks future tactical trends into past folds.",
                "17 Expanding Temporal Folds: Models train exclusively on historical tournaments and test on the next chronological event.",
                "Out-of-Sample Universe: 1,105 strictly out-of-sample evaluated games with zero future leakage.",
                "Scientific Gold Standard: Mirrors real-world deployment where analysts only possess past data before a tournament."
            ], "Random cross-validation is invalid in sports. We engineered 17 chronological walk-forward folds."),

            ("Section 5: Supervised Modeling & ML", "Supervised Model Benchmark Results", "4-Tier Comparison Across 1,105 Out-of-Sample Matches", [
                "Classification Benchmark (Win/Loss): Naive (Brier 0.250) -> Logistic (0.210) -> ElasticNet (0.208) -> LightGBM (0.1967).",
                "LightGBM Classification Performance: Brier = 0.1967 | LogLoss = 0.5741 | AUC-ROC = 0.7613.",
                "Regression Benchmark (Point Margin): Naive (MAE 14.82) -> Ridge (12.35) -> ElasticNet (12.18) -> LightGBM (11.74 pts).",
                "Epistemological Takeaway: The model extracts real predictive signal from pre-game features under strict temporal testing."
            ], "LightGBM outperformed linear baselines across all 1,105 out-of-sample matches."),

            ("Section 5: Supervised Modeling & ML", "Probability Calibration & Expected Calibration Error", "Why Calibrated Probabilities Matter to a Coaching Staff", [
                "The Concept: A 70% win probability model should see exactly 70 wins out of 100 historical predictions.",
                "Empirical Result: LightGBM achieves an Expected Calibration Error of ECE = 0.0314 (3.14% deviation).",
                "Coach-Facing Utility: Win probabilities can be trusted as reliable frequencies rather than uncalibrated confidence scores.",
                "Reliability Diagram: Confirms monotonic alignment across decile bins without overconfident tail distortion."
            ], "Our model achieves an ECE of 0.0314, meaning predicted probabilities match real historical win frequencies."),

            ("Section 5: Supervised Modeling & ML", "Non-Parametric Statistical Inference", "Clustered Bootstrap and Permutation Testing", [
                "Clustered Bootstrap (B = 5,000 Iterations): Computes robust 95% CIs accounting for intra-tournament team correlation.",
                "Permutation Hypothesis Testing (P = 10,000 Shuffles): Validates tactical differentials without assuming normality.",
                "Multiple Testing Control: Benjamini-Hochberg False Discovery Rate controlled strictly at Q = 0.05.",
                "Uncertainty Transparency: Every metric is presented with explicit sample bounds rather than a misleading single point."
            ], "We use clustered bootstrap and permutation tests to report robust confidence intervals."),

            # SECTION 6
            ("Section 6: Feature Attribution", "What Drives Model Predictions?", "Feature Attribution Hierarchy (Permutation & SHAP)", [
                "1. Multi-Tournament Net Rating Differential: Primary predictor of baseline team strength (+0.048 importance).",
                "2. Effective Field Goal Disparity (eFG%): Shooting efficiency differential (+0.035 importance).",
                "3. Turnover Percentage Differential (TOV%): Possession security and transition prevention (+0.024 importance).",
                "4. Offensive Rebounding Disparity (ORB%): Second-chance generation (+0.018 importance).",
                "5. In-Tournament Momentum: Recent group stage margin trajectory (+0.012 importance)."
            ], "Historical Net Rating, shooting efficiency, and turnover control are the primary drivers of model probability."),

            ("Section 6: Feature Attribution", "The Golden Rule: Feature Importance != Causality", "Maintaining Scientific and Epistemological Humility", [
                "Association vs Intervention: Finding that turnover margin is predictive does not prove telling players not to pass causes wins.",
                "Omitted Variable Bias: Latent talent, tactical adjustments, and injuries drive both stats and victories.",
                "The Analyst's Duty: Clearly communicate to coaches that statistical features identify risk patterns, not guaranteed levers.",
                "Professional Framing: Use model attributions to ask better tactical questions, not dictate tactical dogma."
            ], "Feature importance is not causality. Statistical models identify associations; coaches design interventions."),

            # SECTION 7
            ("Section 7: Tournament Simulation", "Why Game-Level Probabilities Are Insufficient", "Propagating Uncertainty Through Dependent Tournament Brackets", [
                "The Limitation: Single-game probabilities cannot evaluate multi-stage tournament advancement paths.",
                "Knockout Volatility: Single-elimination formats compound variance; one cold shooting night causes elimination.",
                "Monte Carlo Propagation: Simulating full tournament brackets from group stages through the championship final.",
                "Output Metrics: P(Advance Group), P(Reach QF), P(Reach SF), P(Reach Final), and P(Win Championship)."
            ], "Single-game win probabilities cannot answer tournament questions. We simulate the entire bracket."),

            ("Section 7: Tournament Simulation", "Retrospective Historical Simulation Findings", "180,000 Monte Carlo Iterations Across 18 Tournaments", [
                "Champion Top-1 Hit Rate: 72.2% (13 of 18 actual champions were simulated #1 pre-tournament favorites).",
                "Champion Top-2 Hit Rate: 77.8% (14 of 18 champions ranked in Top 2).",
                "Champion Top-4 Hit Rate: 100.0% (18 of 18 champions ranked in Top 4).",
                "Mean Champion Rank: 1.50 across 20 years of international basketball.",
                "Sample Size Caveat: N = 18 is a small retrospective sample; metrics prove historical consistency, not clairvoyance."
            ], "Across 180,000 simulations, 100% of historical champions ranked in the model's Top 4."),

            ("Section 7: Tournament Simulation", "Probability Shrinkage & Scenario Sensitivity", "Testing Robustness to Game-Level Probability Overconfidence", [
                "Shrinkage Transformation: p_shrunk = lambda * p + (1 - lambda) * 0.50 across lambda in {0.50, 0.75, 1.00}.",
                "Decision Invariance: Champion Top-1 Hit Rate (72.2%) and Mean Rank (1.50) remain 100% identical across all shrinkage levels.",
                "Uncertainty Compression: Title probabilities shrink naturally (55.0% -> 47.1%) while relative contender rankings hold firm.",
                "Executive Takeaway: Decision support rankings do not depend on fragile, overconfident probability tails."
            ], "Even when shrinking probabilities by 50%, team rankings and champion identification remain perfectly stable."),

            ("Section 7: Tournament Simulation", "Controlled Flagship Counterfactuals", "Replaying Historic Knockout Scenarios (10,000 Runs Each)", [
                "Beijing 2008 Final Replay (Spain vs USA): Under pre-game odds (P=26.4%), Spain won gold in 26.84% of 10k replays.",
                "EuroBasket 2015 Spain Pre-Knockout Path: Despite 2 group losses, Spain retained a 67.6% model-implied title probability.",
                "EuroBasket 2022 Tactical Perturbation: Spain's title odds compressed from 72.0% to 66.2% under lambda=0.75.",
                "Simulation Utility: Quantifies the true mathematical volatility of single-elimination tournament basketball."
            ], "We used simulation to replay historical what-ifs, proving that Spain had a 26.8% chance to beat the Redeem Team in 2008."),

            # SECTION 8
            ("Section 8: Analyst Decision System", "MVP-8 Decision System Architecture", "Synthesizing 6 Heterogeneous Layers into Auditable Dossiers", [
                "Multi-Criteria Score: S_rec = 0.25*Role + 0.25*TS% + 0.20*Reliability + 0.15*PredictiveImpact + 0.15*FilmQuality.",
                "Confidence Tiers: Tier A (High Confidence, N>=150m + Film IRR) | Tier B (Moderate, N>=90m) | Tier C (Limited).",
                "Recommendation Status: RECOMMENDED (Score >= 70.0) | PROCEED WITH CAUTION (55-70) | NOT RECOMMENDED (<55).",
                "Contradiction Audit: Explicitly searches for and highlights discrepancies between stats and film."
            ], "MVP-8 synthesizes stats, roles, film, ML predictions, and simulations into structured candidate dossiers."),

            ("Section 8: Analyst Decision System", "Case Study 1: Lorenzo Brown (EuroBasket 2022)", "Naturalized Guard Integration During Generational Rebuild", [
                "Context: Spain faced an acute backcourt deficit following the Gasol retirements and Ricky Rubio's ACL tear.",
                "Dossier Metrics: Role Fit = 100.0 (Primary Initiator) | TS% = 58.4% | Predictive Net Impact = +4.2 | Tier B Confidence.",
                "Final Recommendation Score: 84.9 (RECOMMENDED).",
                "Historical Outcome: Lorenzo Brown made the All-Tournament Team (15.2 PPG, 7.6 APG) and led Spain to Gold."
            ], "In 2022, the system gave Lorenzo Brown an 84.9 recommendation score, and he led Spain to an unexpected Gold Medal."),

            ("Section 8: Analyst Decision System", "Case Study 2: Pau Gasol (EuroBasket 2015)", "Interior Hub Dominance and Single-Tournament MVP Peak", [
                "Context: Spain lacked exterior creators (Navarro/Rubio absent) and required an offense centered on post hub gravity.",
                "Dossier Metrics: Role Fit = 85.0 | TS% = 64.8% (25.6 PPG, 8.8 RPG) | Video Film Quality = 3.75/4.0 | Tier A Confidence.",
                "Final Recommendation Score: 80.8 (RECOMMENDED).",
                "Historical Outcome: Gasol won Tournament MVP (40 pts vs France in Semis) and led Spain to Gold."
            ], "In 2015, the system awarded Pau Gasol an 80.8 score with Tier A confidence, leading to a legendary MVP run."),

            ("Section 8: Analyst Decision System", "Case Study 3: Ricky Rubio (FIBA World Cup 2019)", "Transition from Pure Distributor to Tournament MVP Scorer", [
                "Context: Spain needed perimeter scoring alongside Marc Gasol's high-post facilitation.",
                "Dossier Metrics: Role Fit = 95.0 (Primary Initiator) | TS% = 56.2% (16.4 PPG, 6.0 APG) | Predictive Impact = +3.8.",
                "Final Recommendation Score: 72.2 (RECOMMENDED).",
                "Historical Outcome: Ricky Rubio won World Cup MVP and led Spain to the World Championship in Beijing."
            ], "In 2019, the system recognized Rubio's role fit and awarded a 72.2 score, outranking naive PPG rules."),

            ("Section 8: Analyst Decision System", "Case Study 4: Calderón vs Navarro (EuroBasket 2011)", "Hyper-Efficiency vs Historic Volume Shot Creation", [
                "Jose Manuel Calderon: TS% = 62.5%, AST/TOV = 4.2, Recommendation Score = 74.8 (High efficiency, ball security).",
                "Juan Carlos Navarro: TS% = 54.8%, USG% = 28.5%, Recommendation Score = 71.9 (High-volume creator, difficult shotmaker).",
                "Analyst Trade-Off: System highlights Calderon as the safer floor general, while validating Navarro as a primary creator.",
                "Coaching Synergy: Coach Scariolo paired both in tandem to secure the 2011 European Championship."
            ], "In 2011, the system highlighted Calderon's efficiency alongside Navarro's volume, and both won Gold together."),

            # SECTION 9
            ("Section 9: Contradiction Surfacing", "Surfacing Contradictions Rather Than Hiding Them", "The Mark of Professional Decision Support", [
                "Small Sample Traps: High scoring efficiency in under 40 minutes is flagged as 'Tier C Insufficient Sample'.",
                "Tactical Mismatches: Strong offensive production paired with blown P&R drop coverage triggers a 'Tactical Contradiction Alert'.",
                "Spacing Deficits: Identifying lineup role overlap when multiple ball-dominant guards share the court.",
                "The Analyst's Duty: Protect the coaching staff from statistical blind spots and small-sample illusions."
            ], "A professional analyst highlights contradictions and small-sample risks rather than hiding them."),

            ("Section 9: Contradiction Surfacing", "Historical Decision Validation vs Baseline Rules", "80.0% Exact Historical Concordance Across Flagship Scenarios", [
                "MVP-8 Multi-Layer Decision System: 80.0% Concordance (4 of 5 Decisions) | 100% Contender Capture.",
                "Baseline Rule A (Naive PPG): 60.0% Concordance (Vulnerable to inefficient high-volume scoring).",
                "Baseline Rule B (Historical Experience): 60.0% Concordance (Vulnerable to aging and role decline).",
                "Scientific Integrity: 5 reconstructed decisions demonstrate qualitative consistency, not proof of causal superiority."
            ], "Our multi-layer decision system outperformed naive counting stat baselines across historical scenarios."),

            ("Section 9: Contradiction Surfacing", "Translation to Basketball Practice", "What Coaching Staff and Sporting Directors Actually Receive", [
                "What a Head Coach Receives: Opponent offensive scheme profiles, P&R coverage vulnerabilities, and film study questions.",
                "What an Assistant Coach Receives: Lineup spacing ratings, shooting variance bounds, and matchup defensive assignments.",
                "What a Sporting Director Receives: Roster archetype balance audits, generational transition profiles, and medal simulations.",
                "Executive Principle: Translation from complex data science into actionable, intuitive basketball language."
            ], "We translate complex data science into clear, tactical language tailored for coaches and directors."),

            # SECTION 10
            ("Section 10: Boundaries & Workflow", "What the System Does NOT Do", "Explicit Boundaries and Epistemological Humility", [
                "The System Does NOT: Replace coaches or scouts, claim causal certainty, or guarantee tournament victory.",
                "The System Does NOT: Provide live in-game tracking or execute transfer-market club scouting.",
                "The System DOES: Organize multi-modal evidence, eliminate future data leakage, and quantify uncertainty.",
                "The System DOES: Deliver calibrated probabilities (ECE = 0.0314) and surface tactical contradictions transparently."
            ], "This slide explicitly defines our boundaries: we empower coaches with evidence, not replace them."),

            ("Section 10: Boundaries & Workflow", "Professional Analyst Operational Workflow", "Implemented Pre-Game Layer vs Potential Live Applications", [
                "Implemented Pre-Game Layer: Multi-tournament historical baselines, archetype discovery, calibrated ML, and bracket simulations.",
                "Potential In-Tournament Workflow: Updating Bayesian priors after group stage games, tracking performance deviations.",
                "Potential Post-Tournament Workflow: Tactical audit, rotation efficiency evaluation, and generational succession planning.",
                "Operational Distinction: Clear separation between certified historical capabilities and live operational extensions."
            ], "We clearly distinguish between our certified historical implementation and potential live workflows."),

            ("Section 10: Boundaries & Workflow", "Technical Stack & Software Engineering Rigor", "Production-Grade Code Quality and 100% Test Automation", [
                "Core Technologies: Python 3.14, DuckDB, Pandas, NumPy, Scikit-Learn, LightGBM, Matplotlib.",
                "Automated Test Suite: 128 Automated Unit & Integration Tests running in 31.8 seconds (100% Pass Rate).",
                "Bitwise Determinism: Master random seed 42 ensures identical simulation dataframes and reproducible results.",
                "Warehouse Rigor: Relational schema constraints, foreign key validation, and cryptographic SHA-256 provenance."
            ], "Our entire repository is backed by 128 automated tests passing 100% in under 32 seconds."),

            ("Section 10: Boundaries & Workflow", "Data Leakage Prevention Architecture", "Methodological Superiority in Sports Time-Series", [
                "Temporal Isolation Barrier: All feature engineering strictly isolates historical past from future evaluation games.",
                "No Duplicate Match Representation: Bilateral perspective is preserved without polluting train/test splits.",
                "Retrospective Outcome Isolation: Tournament champions and medalists are isolated as retrospective targets.",
                "Audit Trail: Complete logging of fold manifests, random seeds, and feature transformations."
            ], "Our strict temporal isolation barriers guarantee zero future data leakage."),

            ("Section 10: Boundaries & Workflow", "Transparent Limitations", "10 Explicit Constraints Acknowledged", [
                "1. Small Tournament Sample Size (N = 18 Tournaments).",
                "2. Small Within-Tournament Game Samples (5 to 9 games per team).",
                "3. Roster Turnover between international FIBA windows.",
                "4. Absence of optical tracking coordinates or live biometric telemetry.",
                "5. Qualitative video coding limited to 420 high-leverage possessions.",
                "6. Non-linear tactical dependencies are difficult to fully parameterize.",
                "7. Historical tournament formats evolved across the 20-year sample.",
                "8. Simulation accuracy is strictly conditional on underlying model calibration.",
                "9. Historical decision validation sample is small (N = 5 Reconstructed Scenarios).",
                "10. Feature attribution represents statistical association, not causal intervention."
            ], "A senior analyst makes uncertainty visible. Here we document 10 explicit limitations."),

            # SECTION 11
            ("Section 11: Core Value of Analyst", "The Value Chain of Basketball Analytics", "Where the Analyst Sits in the Decision-Making Process", [
                "Raw Data -> Structured Information -> Calibrated Evidence -> Tactical Context -> Quantified Uncertainty.",
                "Decision Support Dossier -> Human Expert Judgment -> Final Basketball Decision (Head Coach / Director).",
                "The Analyst's Contribution: Transforming chaos into structured evidence and surfacing actionable trade-offs.",
                "The Coach's Authority: Exercising tactical judgment, emotional intelligence, and in-game leadership."
            ], "This diagram captures our core philosophy: the analyst provides structured evidence; the coach decides."),

            ("Section 11: Core Value of Analyst", "Conclusion & Professional Portfolio Summary", "From Data to Evidence. From Evidence to Better Decisions.", [
                "Summary: Built a complete, certified 20-year decision support system for international basketball.",
                "Demonstrated Competencies: Data Engineering, Econometrics, Supervised ML, Monte Carlo Simulation, Decision Science.",
                "Final Message: 'This project is not an attempt to automate basketball. It is a demonstration of how a dedicated analyst elevates the decision process.'",
                "From Data to Evidence. From Evidence to Better Decisions. Thank you."
            ], "Thank you. This project proves how a data analyst elevates basketball decision-making."),

            # SECTION 12
            ("Section 12: Interview Q&A Appendix", "Coaching & Sporting Leadership Q&A", "Model Answers for Basketball Professionals", [
                "Q: 'How would I use this during tournament preparation?' -> Baseline opponent tendencies, creator gravity, and P&R drop vulnerabilities.",
                "Q: 'What happens when the data disagrees with my video scout?' -> The system surfaces the contradiction as a priority discussion topic.",
                "Q: 'Can I trust a 75% win probability?' -> Yes, because the model is calibrated (ECE = 0.0314); 75% means winning 3 of 4 and losing 1 of 4.",
                "Q: 'What decisions does this support?' -> Roster composition, role balance, lineup spacing, and tournament scenario planning."
            ], "Here are concise model answers to common questions from coaching and sporting leadership."),

            ("Section 12: Interview Q&A Appendix", "Technical Data Science & Engineering Q&A", "Model Answers for Quantitative Interviewers", [
                "Q: 'Why temporal walk-forward validation?' -> K-fold leaks future information; walk-forward preserves chronological realism.",
                "Q: 'Why use Brier score over Accuracy?' -> Accuracy ignores confidence; Brier score evaluates probability calibration directly.",
                "Q: 'Why isn't feature importance causal?' -> Permutation importance measures associative loss degradation, not intervention effects.",
                "Q: 'How do you guarantee reproducibility?' -> Master seed 42, deterministic DuckDB marts, and 128 automated tests."
            ], "Here are technical answers on validation, calibration, causality, and software engineering.")
        ]

        print(f"Building {len(slides_data)}-slide presentation...")
        for i, (sec, title, sub, bullets, notes) in enumerate(slides_data):
            self.add_slide(sec, title, sub, bullets, notes)
            print(f"  Added Slide {i+1:02d}: {title[:40]}...")

        self.prs.save(str(self.output_path))
        print(f"\nSuccessfully generated presentation: {self.output_path} ({len(self.prs.slides)} slides)")


def main():
    builder = PresentationBuilder()
    builder.build_full_presentation()


if __name__ == "__main__":
    main()
