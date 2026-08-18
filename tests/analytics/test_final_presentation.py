"""Automated Tests for the Master Portfolio Presentation Package.

Validates:
1. PowerPoint presentation (.pptx) exists and contains exactly 30 slides.
2. Presentation outline document exists and outlines all 30 slides.
3. Speaker notes document exists and provides oral guidance for all 30 slides.
4. Presentation package README exists.
5. Presentation audit report exists and approves the deck.
6. Zero prohibited claims in presentation documentation.
"""

from pathlib import Path
import pytest
from pptx import Presentation

from src.config import PROJECT_ROOT, REPORTS_DIR


@pytest.fixture(scope="module")
def presentation_dir() -> Path:
    """Fixture providing presentation directory."""
    return PROJECT_ROOT / "portfolio" / "presentation"


@pytest.fixture(scope="module")
def reports_dir() -> Path:
    """Fixture providing reports directory."""
    return REPORTS_DIR


def test_pptx_exists_and_has_30_slides(presentation_dir: Path):
    """Verify PowerPoint file exists and contains exactly 30 slides."""
    pptx_path = presentation_dir / "International_Basketball_Analytics_Presentation.pptx"
    assert pptx_path.exists(), "Presentation .pptx file missing"
    
    prs = Presentation(str(pptx_path))
    assert len(prs.slides) == 30, f"Expected 30 slides, found {len(prs.slides)}"


def test_presentation_outline_covers_30_slides(presentation_dir: Path):
    """Verify presentation_outline.md exists and covers all 30 slides."""
    outline_path = presentation_dir / "presentation_outline.md"
    assert outline_path.exists(), "presentation_outline.md missing"
    
    content = outline_path.read_text(encoding="utf-8")
    for i in range(1, 31):
        slide_str = f"| **{i:02d}**"
        assert slide_str in content, f"Slide {i} missing from outline"


def test_speaker_notes_cover_30_slides(presentation_dir: Path):
    """Verify speaker_notes.md covers all 30 slides."""
    notes_path = presentation_dir / "speaker_notes.md"
    assert notes_path.exists(), "speaker_notes.md missing"
    
    content = notes_path.read_text(encoding="utf-8")
    for i in range(1, 31):
        assert f"Diapositiva {i}" in content, f"Slide {i} missing in speaker notes"


def test_presentation_readme_exists(presentation_dir: Path):
    """Verify presentation README.md exists and has structure."""
    readme_path = presentation_dir / "README.md"
    assert readme_path.exists(), "portfolio/presentation/README.md missing"
    content = readme_path.read_text(encoding="utf-8")
    assert "International_Basketball_Analytics_Presentation.pptx" in content
    assert "30 Diapositivas" in content


def test_presentation_audit_report_exists(reports_dir: Path):
    """Verify final_presentation_audit.md exists and certifies approval."""
    audit_path = reports_dir / "final_presentation_audit.md"
    assert audit_path.exists(), "reports/final_presentation_audit.md missing"
    content = audit_path.read_text(encoding="utf-8")
    assert "APROBADO SIN RESERVAS" in content
    assert "18 Torneos" in content
    assert "1.145" in content


def test_no_banned_language_in_presentation_docs(presentation_dir: Path):
    """Verify absence of prohibited overclaiming words in presentation docs."""
    banned_phrases = [
        "ai coach", "replaces the coach", "guaranteed wins",
        "predicts the winner with 100% accuracy", "adivina el ganador"
    ]
    check_files = [
        presentation_dir / "presentation_outline.md",
        presentation_dir / "speaker_notes.md",
        presentation_dir / "README.md"
    ]
    for p in check_files:
        if p.exists():
            content = p.read_text(encoding="utf-8").lower()
            for b in banned_phrases:
                assert b not in content, f"Banned phrase '{b}' found in {p.name}"
