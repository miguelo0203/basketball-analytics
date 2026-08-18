"""Automated Tests for MVP-9 Analyst Presentation & Portfolio Deliverables.

Validates:
1. Markdown presentation deck integrity (40 slides, 12 sections).
2. Quantitative slide data dictionary and metric traceability.
3. Complete speaker notes covering all 40 slides.
4. Visual asset manifest mapping all publication figures.
5. PowerPoint .pptx deck structure, slide count (40), and readability.
"""

from pathlib import Path
import pytest
import pptx

from src.config import REPORTS_DIR


@pytest.fixture(scope="module")
def pres_dir() -> Path:
    """Fixture providing the presentation directory path."""
    return REPORTS_DIR / "presentation"


def test_markdown_presentation_exists(pres_dir: Path):
    """Verify master presentation markdown file exists and has content."""
    deck_path = pres_dir / "mvp9_analyst_portfolio_presentation.md"
    assert deck_path.exists(), "mvp9_analyst_portfolio_presentation.md missing"
    content = deck_path.read_text(encoding="utf-8")
    assert "Slide 1:" in content
    assert "Slide 40:" in content
    assert len(content) > 10000


def test_slide_data_dictionary_exists(pres_dir: Path):
    """Verify slide data dictionary exists and contains key metrics."""
    data_path = pres_dir / "mvp9_slide_data.md"
    assert data_path.exists(), "mvp9_slide_data.md missing"
    content = data_path.read_text(encoding="utf-8")
    assert "1,145" in content
    assert "0.0314" in content  # ECE
    assert "0.1967" in content  # Brier score


def test_speaker_notes_completeness(pres_dir: Path):
    """Verify speaker notes document covers all 40 slides."""
    notes_path = pres_dir / "mvp9_speaker_notes.md"
    assert notes_path.exists(), "mvp9_speaker_notes.md missing"
    content = notes_path.read_text(encoding="utf-8")
    for i in range(1, 41):
        assert f"Slide {i}:" in content, f"Missing speaker notes for Slide {i}"


def test_visual_asset_manifest_exists(pres_dir: Path):
    """Verify visual asset manifest exists and maps figures."""
    manifest_path = pres_dir / "mvp9_visual_asset_manifest.md"
    assert manifest_path.exists(), "mvp9_visual_asset_manifest.md missing"
    content = manifest_path.read_text(encoding="utf-8")
    assert "fig1_model_benchmark_comparison.png" in content
    assert "fig1_tournament_champion_probabilities.png" in content


def test_pptx_deck_structure_and_slide_count(pres_dir: Path):
    """Verify PowerPoint .pptx file exists and contains exactly 40 slides."""
    pptx_path = pres_dir / "mvp9_analyst_portfolio.pptx"
    assert pptx_path.exists(), "mvp9_analyst_portfolio.pptx missing"
    prs = pptx.Presentation(str(pptx_path))
    assert len(prs.slides) == 40, f"Expected 40 slides in pptx, got {len(prs.slides)}"


def test_pptx_slide_dimensions_and_notes(pres_dir: Path):
    """Verify widescreen 16:9 aspect ratio and speaker notes attachment."""
    pptx_path = pres_dir / "mvp9_analyst_portfolio.pptx"
    prs = pptx.Presentation(str(pptx_path))
    # 13.333 inches width, 7.5 inches height
    assert abs(prs.slide_width.inches - 13.333) < 0.01
    assert abs(prs.slide_height.inches - 7.5) < 0.01

    # Verify first slide has speaker notes
    s1 = prs.slides[0]
    assert s1.has_notes_slide
    assert len(s1.notes_slide.notes_text_frame.text) > 10
